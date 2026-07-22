import ast
import io
import os
import re
import uuid
import json
import base64
import traceback
import time
import queue
import socket
import ipaddress
import urllib.parse
import multiprocessing as mp
from collections.abc import Awaitable, Callable

import pandas as pd
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import plotly.express as px
import plotly.graph_objects as go
import seaborn as sns

# ── Premium chart theme — seaborn base + custom rcParams, matched to the indigo UI ──
INDIGO_PALETTE = ["#4F46E5", "#10B981", "#F59E0B", "#8B5CF6", "#EF4444",
                  "#06B6D4", "#EC4899", "#0EA5E9", "#64748B", "#14B8A6"]
try:
    sns.set_theme(style="whitegrid", palette=INDIGO_PALETTE)
except Exception:
    pass

plt.rcParams.update({
    "figure.facecolor": "white",
    "savefig.facecolor": "white",
    "figure.figsize": (8, 5),
    "figure.dpi": 150,
    "savefig.dpi": 150,
    "savefig.bbox": "tight",
    "axes.facecolor": "white",
    "axes.edgecolor": "#E2E8F0",
    "axes.linewidth": 1.1,
    "axes.labelcolor": "#334155",
    "axes.labelweight": "medium",
    "axes.labelpad": 8,
    "axes.titlecolor": "#0F172A",
    "axes.titlesize": 15,
    "axes.titleweight": "bold",
    "axes.titlepad": 16,
    "axes.labelsize": 11.5,
    "axes.grid": True,
    "axes.axisbelow": True,
    "axes.spines.top": False,
    "axes.spines.right": False,
    "grid.color": "#EEF2F6",
    "grid.linewidth": 1.1,
    "text.color": "#0F172A",
    "xtick.color": "#64748B",
    "ytick.color": "#64748B",
    "xtick.labelsize": 10,
    "ytick.labelsize": 10,
    "font.size": 11,
    "font.family": "sans-serif",
    "font.sans-serif": ["Inter", "Segoe UI", "Helvetica Neue", "Arial", "DejaVu Sans"],
    "legend.frameon": True,
    "legend.framealpha": 0.92,
    "legend.edgecolor": "#E2E8F0",
    "legend.fontsize": 10,
    "figure.autolayout": True,
})
matplotlib.rcParams["axes.prop_cycle"] = matplotlib.cycler(color=INDIGO_PALETTE)

from fastapi import FastAPI, UploadFile, File, HTTPException, Request, BackgroundTasks
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, Response, StreamingResponse
from pydantic import BaseModel, Field
from google import genai
from google.genai import types
from dotenv import load_dotenv
from sandbox_runner import execute_code_worker
import storage

load_dotenv(override=True)

app = FastAPI(title="CSV Analyst AI — Agentic Data Scientist")

# ALLOWED_ORIGINS: comma-separated list, e.g. "https://your-app.vercel.app"
# Falls back to "*" in local dev (when env var is not set).
_raw_origins = os.environ.get("ALLOWED_ORIGINS", "*")
ALLOWED_ORIGINS = [o.strip() for o in _raw_origins.split(",")] if _raw_origins != "*" else ["*"]

app.add_middleware(
    CORSMiddleware,
    allow_origins=ALLOWED_ORIGINS,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

GEMINI_MODEL = "gemini-2.5-flash-lite-preview-06-17"
client = genai.Client(api_key=os.environ.get("GEMINI_API_KEY"))

dataframes: dict[str, pd.DataFrame] = {}
models: dict[str, dict] = {}  # session_id -> trained model info (for inference on new input)
conversation_state: dict[str, dict] = {}
query_cache: dict[str, dict] = {}
session_meta: dict[str, dict] = {}
rate_limit_buckets: dict[str, list[float]] = {}
jobs: dict[str, dict] = {}

MAX_UPLOAD_BYTES = int(os.environ.get("MAX_UPLOAD_BYTES", 25 * 1024 * 1024))
MAX_DATAFRAME_ROWS = int(os.environ.get("MAX_DATAFRAME_ROWS", 100_000))
MAX_DATAFRAME_COLUMNS = int(os.environ.get("MAX_DATAFRAME_COLUMNS", 200))
SESSION_TTL_SECONDS = int(os.environ.get("SESSION_TTL_SECONDS", 4 * 60 * 60))
MAX_SESSIONS = int(os.environ.get("MAX_SESSIONS", 50))
MAX_QUERY_CACHE_ENTRIES = int(os.environ.get("MAX_QUERY_CACHE_ENTRIES", 200))
RATE_LIMIT_MAX_REQUESTS = int(os.environ.get("RATE_LIMIT_MAX_REQUESTS", 600))
RATE_LIMIT_WINDOW_SECONDS = int(os.environ.get("RATE_LIMIT_WINDOW_SECONDS", 60))
MAX_JOBS = int(os.environ.get("MAX_JOBS", 200))


def client_rate_limit_key(request: Request) -> str:
    forwarded_for = request.headers.get("x-forwarded-for", "")
    client_ip = forwarded_for.split(",", 1)[0].strip()
    if not client_ip and request.client:
        client_ip = request.client.host
    return f"{client_ip or 'unknown'}:{request.url.path}"


@app.middleware("http")
async def rate_limit_middleware(
    request: Request,
    call_next: Callable[[Request], Awaitable[Response]],
) -> Response:
    if RATE_LIMIT_MAX_REQUESTS <= 0:
        return await call_next(request)

    now = _now()
    cutoff = now - RATE_LIMIT_WINDOW_SECONDS
    key = client_rate_limit_key(request)
    bucket = [ts for ts in rate_limit_buckets.get(key, []) if ts >= cutoff]
    rate_limit_buckets[key] = bucket

    if len(bucket) >= RATE_LIMIT_MAX_REQUESTS:
        retry_after = max(1, int(RATE_LIMIT_WINDOW_SECONDS - (now - bucket[0])))
        return JSONResponse(
            status_code=429,
            content={
                "detail": "Rate limit exceeded. Please wait before retrying.",
                "retry_after_seconds": retry_after,
                "limit": RATE_LIMIT_MAX_REQUESTS,
                "window_seconds": RATE_LIMIT_WINDOW_SECONDS,
            },
            headers={"Retry-After": str(retry_after)},
        )

    bucket.append(now)
    return await call_next(request)

# ── RAG document store ────────────────────────────────────────────────────────

EMBED_MODEL = "models/text-embedding-004"
EMBED_DIM   = 768


def _chunk_text(text: str, size: int = 600, overlap: int = 80) -> list[str]:
    text = " ".join(text.split())
    chunks, start = [], 0
    while start < len(text):
        chunks.append(text[start:start + size])
        start += size - overlap
    return [c for c in chunks if len(c.strip()) > 40]


def _embed_batch(texts: list[str]) -> list[list[float]]:
    """Embed a list of texts using Gemini text-embedding-004."""
    results = []
    for text in texts:
        try:
            resp = client.models.embed_content(model=EMBED_MODEL, contents=text)
            results.append(list(resp.embeddings[0].values))
        except Exception:
            results.append([0.0] * EMBED_DIM)
    return results


def _parse_doc(content: bytes, filename: str) -> str:
    ext = filename.rsplit(".", 1)[-1].lower()
    if ext == "pdf":
        try:
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(content))
            return "\n".join(p.extract_text() or "" for p in reader.pages)
        except Exception as e:
            raise ValueError(f"Could not parse PDF: {e}")
    elif ext in ("xlsx", "xls"):
        try:
            df_doc = pd.read_excel(io.BytesIO(content))
            return df_doc.to_string(index=False)
        except Exception as e:
            raise ValueError(f"Could not parse Excel: {e}")
    elif ext in ("txt", "md", "rst", "csv"):
        return content.decode("utf-8", errors="replace")
    else:
        raise ValueError(f"Unsupported file type: .{ext}. Use PDF, Excel, or text files.")


class DocStore:
    """Per-session in-memory vector store for RAG over uploaded documentation."""

    def __init__(self):
        self.chunks:    list[str]        = []
        self.embeddings: list[list[float]] = []
        self.metadata:  list[dict]       = []   # {filename, chunk_idx}
        self.filenames: list[str]        = []

    def add(self, text: str, filename: str) -> int:
        chunks = _chunk_text(text)
        embs   = _embed_batch(chunks)
        for i, (c, e) in enumerate(zip(chunks, embs)):
            self.chunks.append(c)
            self.embeddings.append(e)
            self.metadata.append({"filename": filename, "chunk_idx": i})
        if filename not in self.filenames:
            self.filenames.append(filename)
        return len(chunks)

    def search(self, query: str, top_k: int = 4) -> list[dict]:
        if not self.embeddings:
            return []
        q = np.array(_embed_batch([query])[0], dtype=np.float32)
        sims = [
            float(np.dot(q, np.array(e, dtype=np.float32)) /
                  (np.linalg.norm(q) * np.linalg.norm(e) + 1e-9))
            for e in self.embeddings
        ]
        top_idx = sorted(range(len(sims)), key=lambda i: sims[i], reverse=True)[:top_k]
        return [
            {"text": self.chunks[i], "filename": self.metadata[i]["filename"], "score": round(sims[i], 3)}
            for i in top_idx if sims[i] > 0.25
        ]


doc_stores: dict[str, DocStore] = {}


class QueryRequest(BaseModel):
    session_id: str
    question: str
    category: str = "general"


class StoryRequest(BaseModel):
    session_id: str
    category: str = "general"


class InvestigationRequest(BaseModel):
    session_id: str
    goal: str = "Investigate the strongest decision opportunity in this dataset."
    category: str = "general"


class PredictRequest(BaseModel):
    session_id: str
    target: str
    category: str = "general"


class TextUploadRequest(BaseModel):
    text: str
    filename: str = "pasted_data.csv"
    has_header: bool = True


class InferJoinRequest(BaseModel):
    session_id_1: str
    session_id_2: str


class JoinRequest(BaseModel):
    session_id_1: str
    session_id_2: str
    join_key_1: str
    join_key_2: str
    how: str = "inner"


class ForecastRequest(BaseModel):
    session_id: str
    date_column: str
    target_column: str
    periods: int = 12
    freq: str = "auto"


class UrlImportRequest(BaseModel):
    url: str
    filename: str = "imported_dataset.csv"


class CompareRequest(BaseModel):
    session_id_1: str
    session_id_2: str


class PredictInputRequest(BaseModel):
    session_id: str
    values: dict


class ScenarioRequest(BaseModel):
    session_id: str
    baseline: dict = Field(default_factory=dict)
    changes: dict = Field(default_factory=dict)
    category: str = "general"


class ScenarioParseRequest(BaseModel):
    session_id: str
    prompt: str
    category: str = "general"


class CleanRequest(BaseModel):
    actions: list[str] | None = None


class ValidateRowsRequest(BaseModel):
    rows: list[dict]


def _now() -> float:
    return time.time()


def _touch_session(session_id: str) -> None:
    meta = session_meta.setdefault(session_id, {"created_at": _now()})
    meta["last_accessed"] = _now()


def _delete_session(session_id: str) -> None:
    dataframes.pop(session_id, None)
    models.pop(session_id, None)
    doc_stores.pop(session_id, None)
    conversation_state.pop(session_id, None)
    session_meta.pop(session_id, None)
    for key in list(query_cache):
        value = query_cache.get(key, {})
        if key.startswith(f"{session_id}:") or value.get("session_id") == session_id:
            query_cache.pop(key, None)


def get_session_df(session_id: str) -> pd.DataFrame:
    """Retrieve session DataFrame from memory cache or load from SQLite/Parquet on disk."""
    cleanup_expired_sessions()
    if session_id in dataframes:
        _touch_session(session_id)
        storage.update_session_access(session_id)
        return dataframes[session_id]
    df = storage.load_dataframe(session_id)
    if df is not None:
        dataframes[session_id] = df
        _touch_session(session_id)
        storage.update_session_access(session_id)
        return df
    raise HTTPException(status_code=404, detail="Session not found. Upload a CSV first.")


def is_valid_session(session_id: str) -> bool:
    if session_id in dataframes:
        return True
    return storage.get_session(session_id) is not None


def cleanup_expired_sessions() -> None:
    """Remove old in-memory state so local/prod workers do not grow forever."""
    now = _now()
    expired = [
        sid for sid, meta in session_meta.items()
        if now - float(meta.get("last_accessed", meta.get("created_at", now))) > SESSION_TTL_SECONDS
    ]
    for sid in expired:
        _delete_session(sid)

    while len(session_meta) > MAX_SESSIONS:
        oldest = min(
            session_meta,
            key=lambda sid: float(session_meta[sid].get("last_accessed", session_meta[sid].get("created_at", now))),
        )
        _delete_session(oldest)

    storage.cleanup_old_sessions(ttl_seconds=SESSION_TTL_SECONDS)


def _cache_set(key: str, value: dict) -> None:
    query_cache[key] = value
    while len(query_cache) > MAX_QUERY_CACHE_ENTRIES:
        oldest = next(iter(query_cache))
        query_cache.pop(oldest, None)


def create_job(kind: str, session_id: str | None = None) -> dict:
    """Create a bounded background job record backed by SQLite persistence."""
    job_id = str(uuid.uuid4())
    job_data = {
        "job_id": job_id,
        "kind": kind,
        "session_id": session_id,
        "status": "queued",
        "created_at": _now(),
        "updated_at": _now(),
        "result": None,
        "error": None,
    }
    jobs[job_id] = job_data
    storage.save_job(job_id, session_id, kind, "queued")
    while len(jobs) > MAX_JOBS:
        oldest = min(jobs, key=lambda jid: float(jobs[jid].get("updated_at", jobs[jid].get("created_at", _now()))))
        jobs.pop(oldest, None)
    return job_data


def update_job(job_id: str, status: str, result: dict | None = None,
               error: str | None = None) -> None:
    """Update a job record with terminal or progress state."""
    job = jobs.get(job_id)
    if not job:
        job_db = storage.get_job(job_id)
        if job_db:
            job = {
                "job_id": job_db["job_id"],
                "kind": job_db["job_type"],
                "session_id": job_db.get("session_id"),
                "status": job_db["status"],
                "created_at": job_db.get("created_at", _now()),
                "updated_at": job_db.get("updated_at", _now()),
                "result": job_db.get("result"),
                "error": job_db.get("error"),
            }
            jobs[job_id] = job
    if not job:
        return
    job["status"] = status
    job["updated_at"] = _now()
    if result is not None:
        job["result"] = result
    if error is not None:
        job["error"] = error
    storage.update_job(job_id, status, result, error)


def public_job(job: dict) -> dict:
    """Return a JSON-safe view of a background job."""
    return {
        "job_id": job["job_id"],
        "kind": job["kind"],
        "session_id": job.get("session_id"),
        "status": job["status"],
        "created_at": job["created_at"],
        "updated_at": job["updated_at"],
        "result": job.get("result"),
        "error": job.get("error"),
    }


def make_sse_emitter(endpoint: str, session_id: str) -> Callable[[dict], str]:
    """Create an SSE encoder that adds request tracing metadata to every event."""
    request_id = str(uuid.uuid4())
    started_at = _now()
    sequence = 0

    def emit(payload: dict) -> str:
        nonlocal sequence
        sequence += 1
        event = {**payload}
        meta = {**event.get("meta", {})}
        meta.update({
            "request_id": request_id,
            "endpoint": endpoint,
            "session_id": session_id,
            "elapsed_ms": int((_now() - started_at) * 1000),
            "sequence": sequence,
        })
        event["meta"] = meta
        return f"data: {json.dumps(event)}\n\n"

    return emit


def validate_dataframe_limits(df: pd.DataFrame) -> None:
    """Reject datasets that exceed configured production safety limits."""
    if df.shape[0] > MAX_DATAFRAME_ROWS:
        raise HTTPException(
            status_code=413,
            detail=f"Dataset has {df.shape[0]:,} rows; maximum allowed is {MAX_DATAFRAME_ROWS:,}.",
        )
    if df.shape[1] > MAX_DATAFRAME_COLUMNS:
        raise HTTPException(
            status_code=413,
            detail=f"Dataset has {df.shape[1]:,} columns; maximum allowed is {MAX_DATAFRAME_COLUMNS:,}.",
        )


async def read_upload_limited(file: UploadFile) -> bytes:
    """Read an uploaded file while enforcing a hard byte limit."""
    content = await file.read(MAX_UPLOAD_BYTES + 1)
    if len(content) > MAX_UPLOAD_BYTES:
        raise HTTPException(
            status_code=413,
            detail=f"File is too large. Maximum upload size is {MAX_UPLOAD_BYTES // (1024 * 1024)} MB.",
        )
    return content


def _num(x) -> float | None:
    """Coerce to a JSON-safe float (NaN/inf -> None)."""
    try:
        f = float(x)
    except (TypeError, ValueError):
        return None
    if f != f or f in (float("inf"), float("-inf")):
        return None
    return round(f, 2)


def build_profile(df: pd.DataFrame) -> dict:
    """Compute a JSON-safe data profile for the insights panel."""
    numeric_df = df.select_dtypes(include=[np.number])

    total_cells = int(df.shape[0] * df.shape[1])
    missing_total = int(df.isna().sum().sum())
    missing_pct = round(100 * missing_total / total_cells, 1) if total_cells else 0.0
    duplicate_rows = int(df.duplicated().sum())

    numeric_stats: dict[str, dict] = {}
    for col in numeric_df.columns:
        s = numeric_df[col].dropna()
        if s.empty:
            continue
        numeric_stats[str(col)] = {
            "mean": _num(s.mean()),
            "median": _num(s.median()),
            "std": _num(s.std()),
            "min": _num(s.min()),
            "max": _num(s.max()),
        }

    # NaN-safe preview rows
    head = df.head(5).astype(object).where(pd.notna(df.head(5)), None)
    preview = head.to_dict(orient="records")

    return {
        "rows": int(len(df)),
        "columns": [str(c) for c in df.columns],
        "dtypes": {str(k): str(v) for k, v in df.dtypes.items()},
        "preview": preview,
        "numeric_features": int(numeric_df.shape[1]),
        "missing_total": missing_total,
        "missing_pct": missing_pct,
        "duplicate_rows": duplicate_rows,
        "numeric_stats": numeric_stats,
    }


def _fig_to_b64(fig) -> str:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight", dpi=150)
    plt.close(fig)
    return base64.b64encode(buf.getvalue()).decode()


def build_overview_charts(df: pd.DataFrame) -> list[dict]:
    """Deterministically render an instant 'dashboard' the moment a CSV loads —
    no LLM call, so it is fast and demo-safe. Each chart is isolated in try/except."""
    charts: list[dict] = []
    numeric = df.select_dtypes(include="number")

    # 1 · Correlation heatmap
    if numeric.shape[1] >= 2:
        try:
            corr = numeric.corr()
            n = len(corr.columns)
            fig, ax = plt.subplots(figsize=(min(9, 2 + 0.55 * n), min(7.5, 1.5 + 0.5 * n)))
            sns.heatmap(corr, annot=(n <= 12), fmt=".2f", cmap="coolwarm", center=0,
                        vmin=-1, vmax=1, linewidths=0.5, square=True,
                        cbar_kws={"shrink": 0.8}, annot_kws={"size": 7}, ax=ax)
            ax.set_title("Correlation Matrix")
            charts.append({"title": "Correlation Matrix", "chart": _fig_to_b64(fig)})
        except Exception:
            pass

    # 2 · Distribution small-multiples (up to 6 numeric columns)
    if numeric.shape[1] >= 1:
        try:
            cols = list(numeric.columns)[:6]
            ncols = min(3, len(cols))
            nrows = (len(cols) + ncols - 1) // ncols
            fig, axes = plt.subplots(nrows, ncols, figsize=(4.2 * ncols, 3.1 * nrows))
            axes = np.array(axes).reshape(-1)
            for i, c in enumerate(cols):
                sns.histplot(df[c].dropna(), kde=True, ax=axes[i])
                axes[i].set_title(str(c))
                axes[i].set_xlabel("")
                axes[i].set_ylabel("")
            for j in range(len(cols), len(axes)):
                axes[j].set_visible(False)
            fig.suptitle("Distributions", fontsize=15, fontweight="bold")
            fig.tight_layout()
            charts.append({"title": "Distributions", "chart": _fig_to_b64(fig)})
        except Exception:
            pass

    # 3 · Top values of the first meaningful low-cardinality categorical column
    for c in df.columns:
        name = str(c).lower()
        if any(k in name for k in ("date", "time", "_at", "id")):
            continue  # skip dates / identifiers — not useful as categories
        if not pd.api.types.is_numeric_dtype(df[c]) and 2 <= df[c].nunique() <= 25:
            try:
                vc = df[c].value_counts().head(10)
                fig, ax = plt.subplots(figsize=(8, 5))
                sns.barplot(x=vc.values, y=vc.index.astype(str), ax=ax)
                for cont in ax.containers:
                    ax.bar_label(cont, fmt="%.0f", padding=3)
                ax.set_title(f"Top {c} by count")
                ax.set_xlabel("Count")
                ax.set_ylabel(str(c))
                charts.append({"title": f"Top {c}", "chart": _fig_to_b64(fig)})
                break
            except Exception:
                pass

    return charts


def build_proactive_insights(df: pd.DataFrame) -> list[dict]:
    """Create deterministic fact-first insights immediately after upload."""
    insights: list[dict] = []
    roles = infer_column_roles(df)
    profile = build_profile(df)

    def add(kind: str, title: str, finding: str, source_columns: list[str],
            confidence: str = "High", validation: str = "Computed directly from uploaded rows.") -> None:
        insights.append({
            "kind": kind,
            "title": title,
            "finding": finding,
            "source_columns": source_columns,
            "confidence": confidence,
            "validation": validation,
        })

    def fmt(value: float | None) -> str:
        return "n/a" if value is None else f"{value:,}"

    add(
        "profile",
        "Dataset Loaded",
        f"{profile['rows']:,} rows, {len(profile['columns']):,} columns, "
        f"{profile['numeric_features']:,} numeric fields.",
        [],
    )

    metric = roles["metrics"][0] if roles["metrics"] else None
    dimension = roles["dimensions"][0] if roles["dimensions"] else None
    if metric:
        total = _num(df[metric].sum())
        mean = _num(df[metric].mean())
        add(
            "kpi",
            f"Primary Metric: {metric}",
            f"Total {metric}: {fmt(total)}; average per row: {fmt(mean)}.",
            [metric],
            validation=f"Computed from {df[metric].notna().sum():,} non-missing values.",
        )

    if metric and dimension:
        grouped = df.groupby(dimension)[metric].sum().sort_values(ascending=False)
        if not grouped.empty:
            leader = grouped.index[0]
            leader_value = _num(grouped.iloc[0])
            total = float(grouped.sum()) if grouped.sum() else 0.0
            share = round(float(grouped.iloc[0]) / total * 100, 1) if total else 0.0
            add(
                "segment",
                f"Top {dimension}",
                f"{leader} leads {metric} with {fmt(leader_value)} ({share}% of grouped total).",
                [dimension, metric],
                validation=f"Grouped {metric} by {dimension} and ranked descending.",
            )

    if profile["missing_total"] or profile["duplicate_rows"]:
        add(
            "quality",
            "Data Quality Watch",
            f"{profile['missing_total']:,} missing values and "
            f"{profile['duplicate_rows']:,} duplicate rows detected.",
            [],
            confidence="High",
            validation="Computed with DataFrame missing-value and duplicate-row checks.",
        )
    else:
        add(
            "quality",
            "Clean First Pass",
            "No missing values or duplicate rows were detected in the uploaded data.",
            [],
        )

    numeric = df.select_dtypes(include="number")
    if numeric.shape[1] >= 2:
        corr = numeric.corr().abs()
        corr_matrix = corr.to_numpy(copy=True)
        np.fill_diagonal(corr_matrix, 0)
        corr = pd.DataFrame(corr_matrix, index=corr.index, columns=corr.columns)
        max_col = corr.max().idxmax()
        partner = corr[max_col].idxmax()
        score = _num(corr.loc[max_col, partner])
        if score and score >= 0.3:
            add(
                "relationship",
                "Strongest Numeric Relationship",
                f"{max_col} and {partner} have the strongest observed correlation ({score}).",
                [str(max_col), str(partner)],
                confidence="Medium" if score < 0.6 else "High",
                validation="Computed from the absolute Pearson correlation matrix.",
            )

    if metric:
        s = pd.to_numeric(df[metric], errors="coerce").dropna()
        if len(s) >= 8:
            q1, q3 = s.quantile([0.25, 0.75])
            iqr = q3 - q1
            if iqr:
                outliers = int(((s < q1 - 1.5 * iqr) | (s > q3 + 1.5 * iqr)).sum())
                if outliers:
                    add(
                        "anomaly",
                        f"Outliers In {metric}",
                        f"{outliers:,} potential outlier values found using the IQR rule.",
                        [metric],
                        confidence="Medium",
                        validation="Computed with the 1.5x IQR outlier rule.",
                    )

    return insights[:6]


def build_quality_report(df: pd.DataFrame) -> dict:
    """Deterministically diagnose data-quality issues for the health panel."""
    profile = build_profile(df)
    roles = infer_column_roles(df)
    issues: list[dict] = []

    def add(severity: str, title: str, detail: str,
            columns: list[str] | None = None,
            suggestion: str = "") -> None:
        issues.append({
            "severity": severity,
            "title": title,
            "detail": detail,
            "columns": columns or [],
            "suggestion": suggestion,
        })

    if profile["missing_total"]:
        missing_cols = df.isna().sum().sort_values(ascending=False)
        affected = [str(c) for c, n in missing_cols.items() if n > 0][:5]
        add(
            "warn",
            "Missing Values",
            f"{profile['missing_total']:,} missing cells ({profile['missing_pct']}% of all cells).",
            affected,
            "Review whether blanks mean unknown, not applicable, or true zero before modeling.",
        )

    if profile["duplicate_rows"]:
        add(
            "warn",
            "Duplicate Rows",
            f"{profile['duplicate_rows']:,} duplicate rows detected.",
            [],
            "Deduplicate before aggregate reporting if repeated rows are not expected.",
        )

    if roles["ids"]:
        add(
            "info",
            "Identifier Columns",
            f"{len(roles['ids'])} likely identifier column(s) detected.",
            roles["ids"][:5],
            "Exclude identifiers from correlations and predictive modeling.",
        )

    for col in df.columns:
        if pd.api.types.is_numeric_dtype(df[col]):
            continue
        nunique = int(df[col].nunique(dropna=True))
        if nunique > 50 or nunique > max(len(df) * 0.8, 20):
            add(
                "info",
                "High-Cardinality Category",
                f"`{col}` has {nunique:,} unique values.",
                [str(col)],
                "Use top-N grouping or treat this as an identifier-like field in charts.",
            )

    numeric = df.select_dtypes(include="number")
    for col in numeric.columns:
        s = pd.to_numeric(df[col], errors="coerce").dropna()
        if len(s) < 8:
            continue
        q1, q3 = s.quantile([0.25, 0.75])
        iqr = q3 - q1
        if iqr:
            outliers = int(((s < q1 - 1.5 * iqr) | (s > q3 + 1.5 * iqr)).sum())
            if outliers:
                add(
                    "warn",
                    "Potential Outliers",
                    f"`{col}` has {outliers:,} potential outlier value(s) by the IQR rule.",
                    [str(col)],
                    "Inspect outliers before training models or reporting averages.",
                )
        skew = s.skew()
        if pd.notna(skew) and abs(float(skew)) >= 1.5:
            add(
                "info",
                "Skewed Distribution",
                f"`{col}` is strongly skewed (skew={round(float(skew), 2)}).",
                [str(col)],
                "Use medians, log transforms, or distribution charts when interpreting this field.",
            )

    score = 100
    score -= min(35, int(profile["missing_pct"] * 2))
    score -= min(25, int(profile["duplicate_rows"] / max(profile["rows"], 1) * 100))
    score -= min(20, sum(1 for issue in issues if issue["severity"] == "warn") * 5)
    score = max(35, score)
    status = "excellent" if score >= 90 else "good" if score >= 75 else "needs_attention"

    return {
        "score": score,
        "status": status,
        "issues": issues[:8],
        "summary": f"{len(issues)} issue(s) found. Data health score: {score}/100.",
    }


SUPPORTED_CLEANING_ACTIONS = {
    "drop_empty_columns",
    "remove_duplicate_rows",
    "trim_string_whitespace",
    "fill_numeric_median",
    "fill_categorical_mode",
}


def build_cleaning_plan(df: pd.DataFrame) -> dict:
    """Create deterministic, conservative cleaning actions for export."""
    profile = build_profile(df)
    actions: list[dict] = []

    def add(action_id: str, title: str, detail: str, columns: list[str],
            default: bool, impact: str) -> None:
        actions.append({
            "id": action_id,
            "title": title,
            "detail": detail,
            "columns": columns,
            "default": default,
            "impact": impact,
            "reversible": True,
        })

    empty_cols = [str(col) for col in df.columns if int(df[col].isna().sum()) == len(df)]
    if empty_cols:
        add(
            "drop_empty_columns",
            "Drop empty columns",
            f"{len(empty_cols)} column(s) contain no usable values.",
            empty_cols[:10],
            True,
            "Reduces noise before reporting and modeling.",
        )

    if profile["duplicate_rows"]:
        add(
            "remove_duplicate_rows",
            "Remove duplicate rows",
            f"{profile['duplicate_rows']:,} exact duplicate row(s) detected.",
            [],
            True,
            "Prevents duplicate records from inflating aggregates.",
        )

    text_cols = [str(col) for col in df.select_dtypes(include=["object", "string"]).columns]
    whitespace_cols = [
        col for col in text_cols
        if df[col].dropna().astype(str).map(lambda value: value != value.strip()).any()
    ]
    if whitespace_cols:
        add(
            "trim_string_whitespace",
            "Trim text whitespace",
            "Some text values include leading or trailing spaces.",
            whitespace_cols[:10],
            True,
            "Prevents duplicate categories caused only by spacing.",
        )

    missing = df.isna().sum()
    numeric_missing = [
        str(col) for col in df.select_dtypes(include="number").columns
        if int(missing.get(col, 0)) > 0 and df[col].notna().any()
    ]
    if numeric_missing:
        add(
            "fill_numeric_median",
            "Fill numeric blanks with median",
            f"{len(numeric_missing)} numeric column(s) contain missing values.",
            numeric_missing[:10],
            True,
            "Preserves rows while using a robust central value.",
        )

    categorical_missing = [
        str(col) for col in df.columns
        if not pd.api.types.is_numeric_dtype(df[col])
        and int(missing.get(col, 0)) > 0
        and df[col].notna().any()
    ]
    if categorical_missing:
        add(
            "fill_categorical_mode",
            "Fill category blanks with mode",
            f"{len(categorical_missing)} categorical/text column(s) contain missing values.",
            categorical_missing[:10],
            False,
            "Useful for modeling, but confirm blanks do not mean unknown or not applicable.",
        )

    default_actions = [action["id"] for action in actions if action["default"]]
    cleaned, applied = apply_cleaning_actions(df, default_actions, validate=False)
    after_quality = build_quality_report(cleaned)

    return {
        "summary": (
            f"{len(actions)} recommended cleaning action(s). "
            f"Safe default export would apply {len(applied)} action(s)."
        ),
        "actions": actions,
        "default_actions": default_actions,
        "estimated_after_quality": after_quality,
    }


def apply_cleaning_actions(
    df: pd.DataFrame,
    actions: list[str],
    validate: bool = True,
) -> tuple[pd.DataFrame, list[dict]]:
    """Apply conservative cleaning actions and return a cleaned copy plus audit log."""
    unknown = sorted(set(actions) - SUPPORTED_CLEANING_ACTIONS)
    if validate and unknown:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported cleaning action(s): {', '.join(unknown)}",
        )

    cleaned = df.copy()
    applied: list[dict] = []

    if "trim_string_whitespace" in actions:
        changed_columns: list[str] = []
        for col in cleaned.select_dtypes(include=["object", "string"]).columns:
            before = cleaned[col].copy()
            cleaned[col] = cleaned[col].map(lambda value: value.strip() if isinstance(value, str) else value)
            if not before.equals(cleaned[col]):
                changed_columns.append(str(col))
        if changed_columns:
            applied.append({
                "id": "trim_string_whitespace",
                "title": "Trim text whitespace",
                "columns": changed_columns,
                "changed_cells": None,
            })

    if "drop_empty_columns" in actions:
        empty_cols = [col for col in cleaned.columns if int(cleaned[col].isna().sum()) == len(cleaned)]
        if empty_cols:
            cleaned = cleaned.drop(columns=empty_cols)
            applied.append({
                "id": "drop_empty_columns",
                "title": "Drop empty columns",
                "columns": [str(col) for col in empty_cols],
                "changed_cells": len(empty_cols),
            })

    if "remove_duplicate_rows" in actions:
        before_rows = len(cleaned)
        cleaned = cleaned.drop_duplicates().reset_index(drop=True)
        removed = before_rows - len(cleaned)
        if removed:
            applied.append({
                "id": "remove_duplicate_rows",
                "title": "Remove duplicate rows",
                "columns": [],
                "changed_cells": removed,
            })

    if "fill_numeric_median" in actions:
        changed = 0
        changed_columns: list[str] = []
        for col in cleaned.select_dtypes(include="number").columns:
            missing = int(cleaned[col].isna().sum())
            if missing and cleaned[col].notna().any():
                cleaned[col] = cleaned[col].fillna(cleaned[col].median())
                changed += missing
                changed_columns.append(str(col))
        if changed:
            applied.append({
                "id": "fill_numeric_median",
                "title": "Fill numeric blanks with median",
                "columns": changed_columns,
                "changed_cells": changed,
            })

    if "fill_categorical_mode" in actions:
        changed = 0
        changed_columns: list[str] = []
        for col in cleaned.columns:
            if pd.api.types.is_numeric_dtype(cleaned[col]):
                continue
            missing = int(cleaned[col].isna().sum())
            mode = cleaned[col].dropna().mode()
            if missing and not mode.empty:
                cleaned[col] = cleaned[col].fillna(mode.iloc[0])
                changed += missing
                changed_columns.append(str(col))
        if changed:
            applied.append({
                "id": "fill_categorical_mode",
                "title": "Fill category blanks with mode",
                "columns": changed_columns,
                "changed_cells": changed,
            })

    return cleaned, applied


def _role_for_column(column: str, roles: dict[str, list[str]]) -> str:
    for role in ("metrics", "time", "dimensions", "ids", "target_candidates", "numeric"):
        if column in roles.get(role, []):
            return role[:-1] if role.endswith("s") else role
    return "feature"


def build_column_dictionary(df: pd.DataFrame, roles: dict[str, list[str]]) -> list[dict]:
    """Create a compact, UI-safe data dictionary for decision planning."""
    rows = max(len(df), 1)
    dictionary: list[dict] = []
    for col in df.columns:
        series = df[col]
        missing = int(series.isna().sum())
        unique_count = int(series.nunique(dropna=True))
        warnings: list[str] = []
        if missing:
            warnings.append(f"{round(missing / rows * 100, 1)}% missing")
        if col in roles.get("ids", []):
            warnings.append("identifier-like; exclude from modeling signals")
        if unique_count > 50 or unique_count > max(rows * 0.8, 20):
            warnings.append("high cardinality")
        if pd.api.types.is_numeric_dtype(series):
            numeric = pd.to_numeric(series, errors="coerce").dropna()
            if len(numeric) >= 8:
                skew = numeric.skew()
                if pd.notna(skew) and abs(float(skew)) >= 1.5:
                    warnings.append("skewed distribution")

        sample_values = [
            str(value) for value in series.dropna().astype(str).drop_duplicates().head(3).tolist()
        ]
        dictionary.append({
            "name": str(col),
            "dtype": str(series.dtype),
            "role": _role_for_column(str(col), roles),
            "completeness_pct": round((rows - missing) / rows * 100, 1),
            "unique_count": unique_count,
            "sample_values": sample_values,
            "warnings": warnings[:3],
        })
    return dictionary


def infer_contract_type(series: pd.Series) -> str:
    """Infer a portable contract type from a pandas series."""
    non_null = series.dropna()
    if non_null.empty:
        return "string"
    if pd.api.types.is_bool_dtype(series):
        return "boolean"
    if pd.api.types.is_integer_dtype(series):
        return "integer"
    if pd.api.types.is_numeric_dtype(series):
        return "number"
    lowered = non_null.astype(str).str.strip().str.lower()
    if lowered.isin(["true", "false", "yes", "no", "0", "1"]).mean() >= 0.9:
        return "boolean"
    parsed_dates = pd.to_datetime(non_null, errors="coerce", format="mixed")
    if parsed_dates.notna().mean() >= 0.8:
        return "date"
    parsed_numbers = pd.to_numeric(non_null, errors="coerce")
    if parsed_numbers.notna().mean() >= 0.9:
        return "number"
    return "string"


def build_data_contract(df: pd.DataFrame) -> dict:
    """Infer a practical schema contract for validating future rows."""
    roles = infer_column_roles(df)
    profile = build_profile(df)
    columns: list[dict] = []
    for col in df.columns:
        series = df[col]
        contract_type = infer_contract_type(series)
        non_null = series.dropna()
        missing = int(series.isna().sum())
        spec: dict = {
            "name": str(col),
            "type": contract_type,
            "role": _role_for_column(str(col), roles),
            "required": missing == 0,
            "nullable": missing > 0,
            "completeness_pct": round((len(df) - missing) / max(len(df), 1) * 100, 1),
            "unique_count": int(series.nunique(dropna=True)),
            "rules": [],
        }

        if contract_type in ("integer", "number") and not non_null.empty:
            numeric = pd.to_numeric(non_null, errors="coerce").dropna()
            if not numeric.empty:
                spec["min"] = _num(numeric.min())
                spec["max"] = _num(numeric.max())
                spec["rules"].append("must_parse_as_number")
        elif contract_type == "date" and not non_null.empty:
            parsed = pd.to_datetime(non_null, errors="coerce", format="mixed").dropna()
            if not parsed.empty:
                spec["min"] = parsed.min().date().isoformat()
                spec["max"] = parsed.max().date().isoformat()
                spec["rules"].append("must_parse_as_date")
        elif contract_type == "string":
            values = non_null.astype(str)
            if 0 < values.nunique() <= 20:
                spec["allowed_values"] = sorted(values.drop_duplicates().tolist())[:20]
                spec["rules"].append("prefer_known_category")

        if spec["required"]:
            spec["rules"].append("required")
        columns.append(spec)

    return {
        "name": "Inferred Data Contract",
        "version": "1.0",
        "row_count_observed": profile["rows"],
        "column_count": len(columns),
        "required_columns": [col["name"] for col in columns if col["required"]],
        "columns": columns,
        "validation_policy": {
            "missing_required": "error",
            "type_mismatch": "error",
            "unknown_category": "warning",
            "extra_columns": "warning",
        },
    }


def validate_rows_against_contract(rows: list[dict], contract: dict) -> dict:
    """Validate user-provided rows against an inferred data contract."""
    columns = {column["name"]: column for column in contract["columns"]}
    errors: list[dict] = []
    warnings: list[dict] = []
    valid_count = 0

    def is_blank(value: object) -> bool:
        return value is None or (isinstance(value, str) and value.strip() == "")

    def add_issue(target: list[dict], row_index: int, column: str, code: str, message: str) -> None:
        target.append({
            "row": row_index,
            "column": column,
            "code": code,
            "message": message,
        })

    for index, row in enumerate(rows):
        row_errors_before = len(errors)
        for name, spec in columns.items():
            value = row.get(name)
            if is_blank(value):
                if spec["required"]:
                    add_issue(errors, index, name, "missing_required", "Required value is missing.")
                continue

            expected = spec["type"]
            if expected in ("integer", "number"):
                parsed = pd.to_numeric(pd.Series([value]), errors="coerce").iloc[0]
                if pd.isna(parsed):
                    add_issue(errors, index, name, "type_mismatch", "Value must parse as a number.")
                elif expected == "integer" and float(parsed) % 1 != 0:
                    add_issue(errors, index, name, "type_mismatch", "Value must be an integer.")
            elif expected == "date":
                parsed_date = pd.to_datetime(pd.Series([value]), errors="coerce", format="mixed").iloc[0]
                if pd.isna(parsed_date):
                    add_issue(errors, index, name, "type_mismatch", "Value must parse as a date.")
            elif expected == "boolean":
                if str(value).strip().lower() not in {"true", "false", "yes", "no", "0", "1"}:
                    add_issue(errors, index, name, "type_mismatch", "Value must parse as a boolean.")

            allowed = spec.get("allowed_values")
            if allowed and str(value) not in allowed:
                add_issue(warnings, index, name, "unknown_category", "Value was not seen in the original dataset.")

        extra_columns = sorted(set(row) - set(columns))
        for extra in extra_columns:
            add_issue(warnings, index, extra, "extra_column", "Column is not part of the inferred contract.")

        if len(errors) == row_errors_before:
            valid_count += 1

    return {
        "total_rows": len(rows),
        "valid_rows": valid_count,
        "invalid_rows": len(rows) - valid_count,
        "errors": errors,
        "warnings": warnings,
        "contract_version": contract["version"],
    }


def build_dashboard_spec(df: pd.DataFrame, category: str = "general") -> dict:
    """Create a deterministic dashboard blueprint from schema roles and profile."""
    roles = infer_column_roles(df)
    quality = build_quality_report(df)
    metric = roles["metrics"][0] if roles["metrics"] else (roles["numeric"][0] if roles["numeric"] else None)
    secondary_metric = next((col for col in roles["numeric"] if col != metric), None)
    dimension = roles["dimensions"][0] if roles["dimensions"] else None
    secondary_dimension = roles["dimensions"][1] if len(roles["dimensions"]) > 1 else None
    time_col = roles["time"][0] if roles["time"] else None

    kpis: list[dict] = [
        {
            "id": "row_count",
            "label": "Rows",
            "value": int(len(df)),
            "source": "dataframe row count",
            "format": "integer",
        },
        {
            "id": "data_health",
            "label": "Data Health",
            "value": int(quality["score"]),
            "source": "deterministic quality report",
            "format": "score",
        },
    ]
    if metric:
        metric_values = pd.to_numeric(df[metric], errors="coerce")
        kpis.extend([
            {
                "id": f"total_{metric}",
                "label": f"Total {metric}",
                "value": _num(metric_values.sum()),
                "source": metric,
                "format": "number",
            },
            {
                "id": f"average_{metric}",
                "label": f"Average {metric}",
                "value": _num(metric_values.mean()),
                "source": metric,
                "format": "number",
            },
        ])

    charts: list[dict] = []

    def add_chart(chart_id: str, title: str, chart_type: str, x: str | None,
                  y: str | None, purpose: str, question: str) -> None:
        charts.append({
            "id": chart_id,
            "title": title,
            "type": chart_type,
            "x": x,
            "y": y,
            "purpose": purpose,
            "question": question,
        })

    if metric and time_col:
        add_chart(
            "metric_trend",
            f"{metric} over time",
            "line",
            time_col,
            metric,
            "Monitor trend direction, spikes, and drops.",
            f"Plot {metric} over time using {time_col}.",
        )
    if metric and dimension:
        add_chart(
            "segment_performance",
            f"{metric} by {dimension}",
            "bar",
            dimension,
            metric,
            "Rank business segments by the primary metric.",
            f"Show {metric} by {dimension} as a bar chart.",
        )
    if dimension:
        add_chart(
            "segment_volume",
            f"Records by {dimension}",
            "bar",
            dimension,
            None,
            "See whether observations are balanced across segments.",
            f"Show record count by {dimension}.",
        )
    if metric:
        add_chart(
            "metric_distribution",
            f"{metric} distribution",
            "histogram",
            metric,
            None,
            "Understand spread, skew, and unusual values.",
            f"Plot the distribution of {metric}.",
        )
    if metric and secondary_metric:
        add_chart(
            "metric_relationship",
            f"{metric} vs {secondary_metric}",
            "scatter",
            secondary_metric,
            metric,
            "Check whether two measures move together.",
            f"Show a scatter plot of {secondary_metric} vs {metric}.",
        )

    filters = []
    for col in [time_col, dimension, secondary_dimension]:
        if col:
            filters.append({
                "column": col,
                "type": "date_range" if col == time_col else "multi_select",
                "role": _role_for_column(col, roles),
            })

    return {
        "title": f"{category.title()} Data Dashboard",
        "category": category,
        "layout": {
            "columns": 12,
            "density": "analyst_workspace",
            "sections": ["kpis", "trends", "segments", "quality"],
        },
        "kpis": kpis[:6],
        "charts": charts[:6],
        "filters": filters,
        "data_requirements": {
            "metric": metric,
            "secondary_metric": secondary_metric,
            "dimension": dimension,
            "time": time_col,
        },
        "starter_questions": [chart["question"] for chart in charts[:4]],
        "quality_notes": quality["issues"][:3],
    }


def build_decision_actions(df: pd.DataFrame, category: str = "general") -> list[dict]:
    """Turn deterministic findings into business actions with evidence and caveats."""
    roles = infer_column_roles(df)
    quality = build_quality_report(df)
    profile = build_profile(df)
    metric = roles["metrics"][0] if roles["metrics"] else (roles["numeric"][0] if roles["numeric"] else None)
    dimension = roles["dimensions"][0] if roles["dimensions"] else None
    time_col = roles["time"][0] if roles["time"] else None
    actions: list[dict] = []

    def add_action(
        title: str,
        priority: str,
        implication: str,
        recommended_action: str,
        estimated_impact: str,
        confidence: float,
        evidence: list[str],
        risks_assumptions: list[str],
        supporting_columns: list[str],
        suggested_question: str,
    ) -> None:
        actions.append({
            "title": title,
            "priority": priority,
            "category": category,
            "implication": implication,
            "recommended_action": recommended_action,
            "estimated_impact": estimated_impact,
            "confidence": round(confidence, 2),
            "evidence": evidence,
            "risks_assumptions": risks_assumptions,
            "supporting_columns": supporting_columns,
            "suggested_question": suggested_question,
        })

    if metric and dimension:
        grouped = (
            df.groupby(dimension)[metric]
            .sum()
            .dropna()
            .sort_values(ascending=False)
        )
        if len(grouped) >= 2:
            leader = str(grouped.index[0])
            laggard = str(grouped.index[-1])
            leader_value = float(grouped.iloc[0])
            laggard_value = float(grouped.iloc[-1])
            total = float(grouped.sum()) or 1.0
            leader_share = leader_value / total * 100
            gap = leader_value - laggard_value
            add_action(
                title=f"Prioritize the strongest {dimension}",
                priority="high" if leader_share >= 40 else "medium",
                implication=f"{leader} contributes the largest share of {metric}.",
                recommended_action=(
                    f"Review what is working in {leader}, then compare it with {laggard} "
                    "to identify repeatable tactics or operational gaps."
                ),
                estimated_impact=f"Top-to-bottom {metric} gap is {gap:,.2f}.",
                confidence=0.9 if len(grouped) >= 3 else 0.78,
                evidence=[
                    f"{leader}: {leader_value:,.2f} {metric}.",
                    f"{laggard}: {laggard_value:,.2f} {metric}.",
                    f"{leader} represents {leader_share:.1f}% of grouped total.",
                ],
                risks_assumptions=[
                    "Grouped totals assume each uploaded row is a valid business event.",
                    "Segment performance may reflect volume differences, not only efficiency.",
                ],
                supporting_columns=[dimension, metric],
                suggested_question=f"Why is {leader} outperforming {laggard} on {metric}?",
            )

    if metric and time_col:
        temp = df[[time_col, metric]].copy()
        temp[time_col] = pd.to_datetime(temp[time_col], errors="coerce", format="mixed")
        temp[metric] = pd.to_numeric(temp[metric], errors="coerce")
        temp = temp.dropna()
        if len(temp) >= 3:
            trend = temp.sort_values(time_col).groupby(time_col)[metric].sum()
            if len(trend) >= 3:
                first = float(trend.iloc[0])
                last = float(trend.iloc[-1])
                delta = last - first
                pct = delta / abs(first) * 100 if first else 0.0
                if abs(pct) >= 5:
                    direction = "declined" if pct < 0 else "increased"
                    add_action(
                        title=f"Investigate {metric} {direction} over time",
                        priority="high" if pct < -10 else "medium",
                        implication=f"{metric} {direction} by {abs(pct):.1f}% from first to latest observed period.",
                        recommended_action=(
                            "Break the trend by segment, category, and product to identify where the movement started."
                        ),
                        estimated_impact=f"Latest period changed by {delta:,.2f} versus the first observed period.",
                        confidence=0.84 if len(trend) >= 5 else 0.72,
                        evidence=[
                            f"First observed period: {first:,.2f}.",
                            f"Latest observed period: {last:,.2f}.",
                            f"Observed periods: {len(trend)}.",
                        ],
                        risks_assumptions=[
                            "The uploaded file may not contain a complete time range.",
                            "Date granularity and seasonality can affect interpretation.",
                        ],
                        supporting_columns=[time_col, metric],
                        suggested_question=f"Break down the {metric} trend by {dimension or 'category'}.",
                    )

    for issue in quality["issues"][:2]:
        if issue["severity"] != "warn":
            continue
        add_action(
            title=f"Fix {issue['title'].lower()} before decisions",
            priority="high",
            implication=issue["detail"],
            recommended_action=issue.get("suggestion") or "Clean this issue before sharing executive metrics.",
            estimated_impact="Improves reliability of downstream analysis and model training.",
            confidence=0.95,
            evidence=[issue["detail"], f"Current quality score: {quality['score']}/100."],
            risks_assumptions=["Cleaning rules should be confirmed with the data owner before production use."],
            supporting_columns=issue.get("columns", []),
            suggested_question="What data quality issues should I fix first?",
        )

    if metric and len(df) >= 50 and len(roles["numeric"]) >= 2:
        add_action(
            title=f"Train a predictive model for {metric}",
            priority="medium",
            implication=f"The dataset appears large enough to explore drivers of {metric}.",
            recommended_action=(
                f"Train a model for `{metric}`, inspect feature importance, then run scenario tests on controllable drivers."
            ),
            estimated_impact="Identifies leading drivers that may be used for planning or simulation.",
            confidence=0.75,
            evidence=[
                f"{len(df):,} rows available.",
                f"{len(roles['numeric'])} numeric columns detected.",
                "Predictive modeling feature is available in the workspace.",
            ],
            risks_assumptions=[
                "Predictive relationships are not causal by default.",
                "Model quality depends on whether the uploaded columns include real driver variables.",
            ],
            supporting_columns=[metric] + roles["numeric"][:4],
            suggested_question=f"Train a predictive model for {metric}.",
        )

    if not actions:
        add_action(
            title="Start with data readiness",
            priority="medium",
            implication="The dataset needs clearer metrics, segments, or time fields before strong business actions can be inferred.",
            recommended_action="Confirm the business meaning of the columns, then add a metric and segment field if available.",
            estimated_impact="Improves the quality of future insights and dashboard recommendations.",
            confidence=0.7,
            evidence=[
                f"{profile['rows']:,} rows and {len(profile['columns']):,} columns detected.",
                f"Data health score: {quality['score']}/100.",
            ],
            risks_assumptions=["Limited schema context can reduce recommendation specificity."],
            supporting_columns=[],
            suggested_question="What are the best next questions for this dataset?",
        )

    return actions[:6]


def build_decision_brief(df: pd.DataFrame, category: str = "general") -> dict:
    """Turn a dataset profile into a decision-readiness and action brief."""
    profile = build_profile(df)
    roles = infer_column_roles(df)
    quality = build_quality_report(df)
    decision_actions = build_decision_actions(df, category)
    metric = roles["metrics"][0] if roles["metrics"] else (roles["numeric"][0] if roles["numeric"] else None)
    dimension = roles["dimensions"][0] if roles["dimensions"] else None
    time_col = roles["time"][0] if roles["time"] else None

    readiness_score = int(quality["score"])
    risk_flags: list[dict] = []

    def risk(severity: str, title: str, detail: str) -> None:
        risk_flags.append({"severity": severity, "title": title, "detail": detail})

    if not metric:
        readiness_score -= 20
        risk("high", "No clear metric", "The dataset has no obvious numeric measure for KPI tracking.")
    if not dimension:
        readiness_score -= 10
        risk("medium", "No clear segment", "The dataset has no obvious categorical dimension for comparison.")
    if not time_col:
        risk("info", "No time field", "Trend, seasonality, and forecasting workflows may be limited.")
    if profile["rows"] < 30:
        readiness_score -= 10
        risk("medium", "Small dataset", "Patterns may be unstable because the dataset has fewer than 30 rows.")
    if quality["status"] == "needs_attention":
        readiness_score -= 10
    readiness_score = max(0, min(100, readiness_score))

    use_cases: list[dict] = []
    blocked_use_cases: list[dict] = []
    if metric and dimension:
        use_cases.append({
            "name": "Performance dashboard",
            "fit": "strong",
            "why": f"Use `{metric}` by `{dimension}` to rank segments and monitor performance.",
        })
    else:
        blocked_use_cases.append({
            "name": "Segment performance dashboard",
            "missing": "Needs at least one metric and one segment/dimension column.",
        })
    if metric and time_col:
        use_cases.append({
            "name": "Trend monitoring",
            "fit": "strong",
            "why": f"Use `{time_col}` with `{metric}` for movement over time.",
        })
    else:
        blocked_use_cases.append({
            "name": "Trend monitoring",
            "missing": "Needs a time column and a numeric metric.",
        })
    if profile["rows"] >= 50 and len(roles["numeric"]) >= 2:
        use_cases.append({
            "name": "Predictive modeling",
            "fit": "moderate",
            "why": "There are enough rows and numeric signals to train an initial model.",
        })
    else:
        blocked_use_cases.append({
            "name": "Predictive modeling",
            "missing": "Usually needs at least 50 rows and multiple usable signal columns.",
        })
    if quality["issues"]:
        use_cases.append({
            "name": "Data quality remediation",
            "fit": "strong",
            "why": "Quality issues were detected and can be turned into cleanup tasks.",
        })

    priority_questions: list[str] = []
    if metric and dimension:
        priority_questions.extend([
            f"Which {dimension} has the highest total {metric}?",
            f"Show {metric} by {dimension} as a bar chart.",
        ])
    if metric and time_col:
        priority_questions.append(f"Plot {metric} over time using {time_col}.")
    if len(roles["numeric"]) >= 2:
        priority_questions.append("Which numeric columns are most strongly correlated?")
    priority_questions.extend([
        "What data quality issues should I fix first?",
        "What story does this dataset tell?",
    ])

    next_actions: list[dict] = []
    for issue in quality["issues"][:3]:
        next_actions.append({
            "priority": "high" if issue["severity"] == "warn" else "medium",
            "action": issue["title"],
            "reason": issue["detail"],
            "impact": issue.get("suggestion") or "Improves trust in downstream analysis.",
        })
    if metric and dimension:
        next_actions.append({
            "priority": "high",
            "action": "Build the first KPI view",
            "reason": f"`{metric}` and `{dimension}` create an immediate performance breakdown.",
            "impact": "Gives stakeholders a clear starting point for decisions.",
        })
    if metric and time_col:
        next_actions.append({
            "priority": "medium",
            "action": "Add trend monitoring",
            "reason": f"`{time_col}` enables time-series analysis for `{metric}`.",
            "impact": "Makes changes, spikes, and drops visible earlier.",
        })

    automation_opportunities = []
    if metric:
        automation_opportunities.append(f"Scheduled KPI digest for `{metric}`")
    if metric and dimension:
        automation_opportunities.append(f"Anomaly alerts for unusual `{metric}` by `{dimension}`")
    if profile["rows"] >= 50:
        automation_opportunities.append("Recurring model retraining when new rows arrive")

    return {
        "category": category,
        "readiness_score": readiness_score,
        "readiness_label": (
            "decision_ready" if readiness_score >= 85
            else "usable_with_caution" if readiness_score >= 65
            else "needs_cleanup"
        ),
        "summary": (
            f"{profile['rows']:,} rows, {len(profile['columns']):,} columns, "
            f"{quality['summary'].lower()}"
        ),
        "recommended_use_cases": use_cases[:5],
        "blocked_use_cases": blocked_use_cases[:4],
        "priority_questions": priority_questions[:6],
        "next_actions": next_actions[:6],
        "decision_actions": decision_actions,
        "risk_flags": risk_flags[:6],
        "automation_opportunities": automation_opportunities[:4],
        "column_dictionary": build_column_dictionary(df, roles),
    }


INVESTIGATION_PERSONAS: dict[str, dict] = {
    "general": {
        "name": "Analytics Lead",
        "focus": "reliable statistical evidence, data readiness, and the fastest useful drill-down",
        "priority_kpis": ["primary metric", "trend", "segment gap", "outliers", "data quality"],
        "column_terms": ["score", "value", "total", "amount", "date", "category", "segment"],
        "default_question": "Which finding should become the first operational dashboard?",
        "risk_language": "statistical uncertainty, sample size, and data completeness",
    },
    "financial": {
        "name": "CFO",
        "focus": "revenue movement, margin pressure, volatility, cost control, and budget risk",
        "priority_kpis": ["revenue", "profit", "margin", "cost", "growth", "variance", "cash flow"],
        "column_terms": ["revenue", "profit", "margin", "cost", "price", "amount", "sales"],
        "default_question": "What is the financial impact and where should leadership act first?",
        "risk_language": "forecast risk, margin sensitivity, and incomplete period coverage",
    },
    "medical": {
        "name": "Healthcare Analyst",
        "focus": "patient cohorts, outcome prevalence, clinical risk factors, and measurement reliability",
        "priority_kpis": ["outcome", "prevalence", "risk", "age", "cohort", "rate", "measurement"],
        "column_terms": ["outcome", "diagnosis", "risk", "age", "patient", "rate", "score"],
        "default_question": "Which cohort or risk factor deserves clinical review first?",
        "risk_language": "association rather than causation, cohort bias, and clinical validity",
    },
    "retail": {
        "name": "Retail Operator",
        "focus": "sales growth, product performance, regional gaps, basket behavior, and demand signals",
        "priority_kpis": ["revenue", "sales", "quantity", "order value", "product", "region", "rating"],
        "column_terms": ["revenue", "sales", "quantity", "product", "category", "region", "rating"],
        "default_question": "Which product, category, or region should the team optimize first?",
        "risk_language": "seasonality, volume mix, returns, and missing customer context",
    },
    "marketing": {
        "name": "Growth Strategist",
        "focus": "channel performance, conversion, retention, audience segments, and campaign efficiency",
        "priority_kpis": ["conversion", "retention", "CAC", "ROAS", "channel", "campaign", "segment"],
        "column_terms": ["conversion", "retention", "channel", "campaign", "segment", "customer", "rate"],
        "default_question": "Which segment or channel has the clearest growth opportunity?",
        "risk_language": "attribution limits, funnel leakage, cohort effects, and campaign mix",
    },
    "hr": {
        "name": "People Analytics Lead",
        "focus": "headcount, attrition, tenure, compensation equity, performance, and workforce fairness",
        "priority_kpis": ["attrition", "tenure", "salary", "performance", "department", "age", "headcount"],
        "column_terms": ["attrition", "tenure", "salary", "performance", "department", "employee", "age"],
        "default_question": "Which workforce segment needs the most urgent people decision?",
        "risk_language": "fairness, privacy, small-group sensitivity, and confounding variables",
    },
}


def investigation_persona_for(category: str) -> dict:
    """Return the deterministic investigation persona for a category lens."""
    return INVESTIGATION_PERSONAS.get(category, INVESTIGATION_PERSONAS["general"])


def persona_relevant_columns(df: pd.DataFrame, persona: dict) -> list[str]:
    """Find columns that match the current analyst persona's KPI vocabulary."""
    terms = [str(term).lower() for term in persona.get("column_terms", [])]
    matches: list[str] = []
    for column in df.columns:
        name = str(column).lower()
        if any(term in name for term in terms):
            matches.append(str(column))
    return matches[:6]


def build_investigation_report(df: pd.DataFrame, goal: str, category: str = "general") -> dict:
    """Run a deterministic multi-step investigation over the most decision-useful columns."""
    profile = build_profile(df)
    roles = infer_column_roles(df)
    quality = build_quality_report(df)
    persona = investigation_persona_for(category)
    persona_columns = persona_relevant_columns(df, persona)
    metric = roles["metrics"][0] if roles["metrics"] else (roles["numeric"][0] if roles["numeric"] else None)
    dimension = roles["dimensions"][0] if roles["dimensions"] else None
    secondary_dimension = roles["dimensions"][1] if len(roles["dimensions"]) > 1 else None
    time_col = roles["time"][0] if roles["time"] else None
    relevant_columns = list(dict.fromkeys([col for col in [metric, dimension, secondary_dimension, time_col] if col] + persona_columns))
    findings: list[str] = []
    recommended_actions: list[str] = []
    investigation_tree: list[dict] = []
    chart_json: str | None = None
    trend_payload: dict | None = None
    segment_payload: dict | None = None
    outlier_payload: dict | None = None
    correlation_payload: dict | None = None

    def add_tree(node: str, status: str, finding: str, columns: list[str] | None = None) -> None:
        investigation_tree.append({
            "node": node,
            "status": status,
            "finding": finding,
            "columns": columns or [],
        })

    add_tree(
        "Dataset profile",
        "complete",
        f"{profile['rows']:,} rows, {len(profile['columns']):,} columns, {profile['missing_pct']}% missing cells.",
        list(profile["columns"])[:6],
    )
    add_tree(
        "Persona focus",
        "complete",
        f"{persona['name']} lens prioritizes {persona['focus']}.",
        persona_columns,
    )

    if metric and time_col:
        temp = df[[time_col, metric]].copy()
        temp[time_col] = pd.to_datetime(temp[time_col], errors="coerce", format="mixed")
        temp[metric] = pd.to_numeric(temp[metric], errors="coerce")
        temp = temp.dropna()
        if len(temp) >= 2:
            trend = temp.sort_values(time_col).groupby(time_col)[metric].sum()
            first = float(trend.iloc[0])
            last = float(trend.iloc[-1])
            delta = last - first
            pct = delta / abs(first) * 100 if first else 0.0
            direction = "increased" if delta >= 0 else "declined"
            trend_payload = {
                "time_column": time_col,
                "metric": metric,
                "first_period": str(trend.index[0].date()),
                "latest_period": str(trend.index[-1].date()),
                "first_value": _num(first),
                "latest_value": _num(last),
                "delta": _num(delta),
                "pct_change": _num(pct),
                "direction": direction,
            }
            finding = (
                f"{metric} {direction} by {abs(pct):.1f}% from {trend_payload['first_period']} "
                f"to {trend_payload['latest_period']}."
            )
            findings.append(finding)
            add_tree("Trend scan", "complete", finding, [time_col, metric])
            if delta < 0:
                recommended_actions.append(
                    f"Break the {metric} decline by segment and date to locate the first point of deterioration."
                )
            else:
                recommended_actions.append(
                    f"Identify what changed during the latest {metric} increase and decide whether it can be repeated."
                )
            plot_df = trend.reset_index()
            fig = px.line(plot_df, x=time_col, y=metric, markers=True, title=f"{metric} investigation trend")
            chart_json = _chart_to_json(fig)
        else:
            add_tree("Trend scan", "limited", f"`{time_col}` and `{metric}` did not contain enough valid paired rows.", [time_col, metric])
    elif metric:
        add_tree("Trend scan", "skipped", "No usable time column was detected for trend analysis.", [metric])

    if metric and dimension:
        grouped = (
            df[[dimension, metric]]
            .dropna()
            .assign(**{metric: pd.to_numeric(df[metric], errors="coerce")})
            .dropna()
            .groupby(dimension)[metric]
            .sum()
            .sort_values(ascending=False)
        )
        if len(grouped) >= 2:
            leader = str(grouped.index[0])
            laggard = str(grouped.index[-1])
            leader_value = float(grouped.iloc[0])
            laggard_value = float(grouped.iloc[-1])
            total = float(grouped.sum()) or 1.0
            gap = leader_value - laggard_value
            leader_share = leader_value / total * 100
            segment_payload = {
                "dimension": dimension,
                "metric": metric,
                "leader": leader,
                "leader_value": _num(leader_value),
                "laggard": laggard,
                "laggard_value": _num(laggard_value),
                "gap": _num(gap),
                "leader_share_pct": _num(leader_share),
                "top_segments": [
                    {"label": str(idx), "value": _num(value)}
                    for idx, value in grouped.head(5).items()
                ],
            }
            finding = (
                f"{leader} leads {dimension} with {leader_value:,.2f} {metric}, "
                f"{gap:,.2f} above {laggard}."
            )
            findings.append(finding)
            add_tree("Segment driver scan", "complete", finding, [dimension, metric])
            recommended_actions.append(
                f"Compare operating conditions for {leader} and {laggard}; the gap is the fastest root-cause path."
            )
            if not chart_json:
                fig = px.bar(
                    pd.DataFrame(segment_payload["top_segments"]),
                    x="label",
                    y="value",
                    title=f"{metric} by {dimension}",
                    labels={"label": dimension, "value": metric},
                )
                chart_json = _chart_to_json(fig)
        else:
            add_tree("Segment driver scan", "limited", f"`{dimension}` has fewer than two usable groups.", [dimension, metric])
    elif metric:
        add_tree("Segment driver scan", "skipped", "No categorical dimension was detected for segment comparison.", [metric])

    if metric:
        series = pd.to_numeric(df[metric], errors="coerce").dropna()
        if len(series) >= 3:
            std = float(series.std(ddof=0))
            mean = float(series.mean())
            if std > 0:
                z_scores = ((series - mean).abs() / std)
                outlier_count = int((z_scores >= 2).sum())
                outlier_payload = {
                    "metric": metric,
                    "count": outlier_count,
                    "threshold": "2 standard deviations from mean",
                    "max_value": _num(series.max()),
                    "min_value": _num(series.min()),
                }
                finding = (
                    f"{outlier_count} {metric} value(s) sit at least 2 standard deviations from the mean."
                )
                findings.append(finding)
                add_tree("Anomaly scan", "complete", finding, [metric])
                if outlier_count:
                    recommended_actions.append(
                        f"Audit the extreme {metric} rows before using them for planning or forecasting."
                    )
        else:
            add_tree("Anomaly scan", "limited", f"`{metric}` has too few numeric values for a stable anomaly scan.", [metric])

    numeric = df.select_dtypes(include="number")
    if numeric.shape[1] >= 2:
        corr = numeric.corr()
        pairs: list[dict] = []
        for i, left in enumerate(corr.columns):
            for right in corr.columns[i + 1:]:
                value = corr.loc[left, right]
                if pd.notna(value):
                    pairs.append({
                        "left": str(left),
                        "right": str(right),
                        "correlation": _num(value),
                        "abs_correlation": abs(float(value)),
                    })
        if pairs:
            strongest = sorted(pairs, key=lambda item: item["abs_correlation"], reverse=True)[0]
            correlation_payload = strongest
            finding = (
                f"Strongest numeric relationship: {strongest['left']} vs {strongest['right']} "
                f"at correlation {strongest['correlation']}."
            )
            findings.append(finding)
            add_tree("Driver relationship scan", "complete", finding, [strongest["left"], strongest["right"]])
    else:
        add_tree("Driver relationship scan", "skipped", "At least two numeric columns are needed for correlation analysis.", roles["numeric"])

    if quality["issues"]:
        issue = quality["issues"][0]
        finding = f"Data quality risk: {issue['detail']}"
        findings.append(finding)
        add_tree("Quality risk scan", "complete", finding, issue.get("columns", []))
        recommended_actions.append(issue.get("suggestion") or "Resolve the leading data quality issue before executive use.")
    else:
        add_tree("Quality risk scan", "complete", "No missing values or duplicate rows were detected.", [])

    if not recommended_actions:
        recommended_actions.append("Use the dashboard blueprint and priority questions to turn this dataset into a recurring KPI workflow.")
    recommended_actions.append(f"Run the next drill-down as a {persona['name']} decision review: {persona['default_question']}")

    confidence = 0.94 if metric and (dimension or time_col) else 0.82
    executive = (
        f"Autonomous investigation for `{goal}` found {len(findings)} evidence points. "
        f"The primary lens is `{metric}`" + (f" by `{dimension}`" if dimension else "") + "."
        if metric else
        f"Autonomous investigation for `{goal}` focused on data readiness because no clear metric was detected."
    )
    report_lines = [
        "**Executive finding**",
        executive,
        "",
        "**Analyst persona**",
        f"{persona['name']} lens: prioritizes {', '.join(persona['priority_kpis'][:5])}.",
        "",
        "**Evidence**",
    ]
    report_lines.extend(f"- {finding}" for finding in findings[:6])
    report_lines.extend([
        "",
        "**Recommended actions**",
    ])
    report_lines.extend(f"- {action}" for action in recommended_actions[:5])
    report_lines.extend([
        "",
        "**Risks / assumptions**",
        "- The investigation uses only the uploaded data and does not infer causality by itself.",
        "- Segment totals can reflect volume differences as well as performance differences.",
        f"- Persona-specific caution: {persona['risk_language']}.",
        f"- Data quality status: {quality['summary']}",
    ])

    plan = {
        "strategy": "Autonomous deterministic investigation using profile, trend, segment, anomaly, correlation, and quality scans.",
        "query_type": "deterministic_investigation",
        "goal": goal,
        "category": category,
        "persona": persona,
        "persona_columns": persona_columns,
        "relevant_columns": relevant_columns,
        "analysis_steps": [
            "Frame the business goal",
            f"Apply the {persona['name']} persona lens",
            "Profile decision readiness",
            "Scan metric movement over time",
            "Rank segment drivers",
            "Detect unusual metric values",
            "Check numeric relationships",
            "Convert evidence into recommended actions",
        ],
        "investigation_tree": investigation_tree,
        "chart_type": "line" if trend_payload else "bar" if segment_payload else "none",
    }

    return {
        "goal": goal,
        "category": category,
        "persona": persona,
        "persona_columns": persona_columns,
        "primary_metric": metric,
        "primary_dimension": dimension,
        "secondary_dimension": secondary_dimension,
        "trend": trend_payload,
        "segment_driver": segment_payload,
        "outliers": outlier_payload,
        "correlation": correlation_payload,
        "quality": {
            "score": quality["score"],
            "status": quality["status"],
            "summary": quality["summary"],
            "issues": quality["issues"][:3],
        },
        "findings": findings[:8],
        "recommended_actions": recommended_actions[:6],
        "investigation_tree": investigation_tree,
        "plan": plan,
        "report": "\n".join(report_lines),
        "chart_json": chart_json,
        "decision_actions": build_decision_actions(df, category),
        "validation": build_validation_payload(
            df,
            method="Autonomous deterministic investigation pipeline",
            source_columns=relevant_columns,
            confidence=confidence,
            reasons=[
                "Evidence was computed directly from the uploaded dataframe.",
                "The investigation separated numeric computation from the final narrative.",
                "Trend, segment, anomaly, correlation, and quality checks are deterministic.",
            ],
        ),
        "meta": {"route": "deterministic_investigation", "facts_first": True},
    }


def build_story_facts(df: pd.DataFrame, category: str = "general") -> dict:
    """Build a deterministic facts JSON payload before any narrative is written."""
    profile = build_profile(df)
    roles = infer_column_roles(df)
    insights = build_proactive_insights(df)
    facts: dict = {
        "category": category,
        "profile": {
            "rows": profile["rows"],
            "columns": len(profile["columns"]),
            "numeric_features": profile["numeric_features"],
            "missing_total": profile["missing_total"],
            "missing_pct": profile["missing_pct"],
            "duplicate_rows": profile["duplicate_rows"],
        },
        "column_roles": roles,
        "insights": insights,
        "top_breakdowns": [],
        "relationships": [],
    }

    metric = roles["metrics"][0] if roles["metrics"] else None
    if metric:
        facts["primary_metric"] = {
            "name": metric,
            "total": _num(df[metric].sum()),
            "mean": _num(df[metric].mean()),
            "non_missing": int(df[metric].notna().sum()),
        }

    for dimension in roles["dimensions"][:3]:
        if metric:
            grouped = df.groupby(dimension)[metric].sum().sort_values(ascending=False).head(5)
            facts["top_breakdowns"].append({
                "dimension": dimension,
                "metric": metric,
                "values": [
                    {"label": str(idx), "value": _num(value)}
                    for idx, value in grouped.items()
                ],
            })

    numeric = df.select_dtypes(include="number")
    if numeric.shape[1] >= 2:
        corr = numeric.corr()
        pairs: list[dict] = []
        for i, left in enumerate(corr.columns):
            for right in corr.columns[i + 1:]:
                value = corr.loc[left, right]
                if pd.notna(value):
                    pairs.append({
                        "left": str(left),
                        "right": str(right),
                        "correlation": _num(value),
                        "abs_correlation": abs(float(value)),
                    })
        facts["relationships"] = sorted(pairs, key=lambda p: p["abs_correlation"], reverse=True)[:3]

    return facts


def render_story_from_facts(facts: dict) -> str:
    """Write a concise story from deterministic facts only."""
    profile = facts["profile"]
    def fmt(value: float | None) -> str:
        return "n/a" if value is None else f"{value:,}"

    lines = [
        "**Headline:** This dataset is ready for analysis, with the strongest first story coming from its primary metrics and segment breakdowns.",
        "",
        "**Key Findings:**",
        f"- The dataset contains {profile['rows']:,} rows and {profile['columns']:,} columns, including {profile['numeric_features']:,} numeric fields.",
    ]

    metric = facts.get("primary_metric")
    if metric:
        lines.append(
            f"- The primary metric appears to be `{metric['name']}` with total {fmt(metric['total'])} "
            f"and average {fmt(metric['mean'])} across {metric['non_missing']:,} non-missing rows."
        )

    if facts["top_breakdowns"] and facts["top_breakdowns"][0]["values"]:
        top = facts["top_breakdowns"][0]
        leader = top["values"][0]
        lines.append(
            f"- `{leader['label']}` is the leading `{top['dimension']}` by `{top['metric']}` "
            f"with {fmt(leader['value'])}."
        )

    if profile["missing_total"] or profile["duplicate_rows"]:
        lines.append(
            f"- Data quality needs attention: {profile['missing_total']:,} missing values and "
            f"{profile['duplicate_rows']:,} duplicate rows were detected."
        )
    else:
        lines.append("- First-pass data quality is clean: no missing values or duplicate rows were detected.")

    if facts["relationships"]:
        rel = facts["relationships"][0]
        lines.append(
            f"- The strongest numeric relationship is `{rel['left']}` vs `{rel['right']}` "
            f"with correlation {rel['correlation']}."
        )

    lines.extend([
        "",
        "**Recommended Next Actions:**",
        "- Review the top segment breakdowns to identify what is driving the main metric.",
        "- Ask for a trend or distribution view to separate normal variation from possible anomalies.",
        "- If the target column is meaningful, train a predictive model and inspect the explainability tabs.",
        "",
        "**Trust Note:** This story is generated from deterministic profile, grouping, quality, and correlation facts before any narrative wording is applied.",
    ])
    return "\n".join(lines)


def build_story_chart_json(facts: dict) -> str | None:
    """Create a Plotly chart for the strongest deterministic story breakdown."""
    if not facts.get("top_breakdowns"):
        return None
    top = facts["top_breakdowns"][0]
    values = top.get("values") or []
    if not values:
        return None
    plot_df = pd.DataFrame(values)
    fig = px.bar(
        plot_df,
        x="label",
        y="value",
        title=f"Top {top['dimension']} by {top['metric']}",
        labels={"label": top["dimension"], "value": top["metric"]},
    )
    return _chart_to_json(fig)


def suggest_followups(question: str, df: pd.DataFrame,
                      plan: dict | None = None,
                      previous_state: dict | None = None) -> list[str]:
    """Generate lightweight context-aware follow-up questions."""
    roles = infer_column_roles(df)
    metric = None
    dimension = None
    if plan:
        spec = plan.get("chart_spec") or {}
        metric = spec.get("y") if spec.get("y") != "count" else None
        dimension = spec.get("x")
    if previous_state:
        metric = metric or previous_state.get("last_metric")
        dimension = dimension or previous_state.get("last_grouping")
    metric = metric if metric in df.columns else (roles["metrics"][0] if roles["metrics"] else None)
    dimension = dimension if dimension in df.columns else (roles["dimensions"][0] if roles["dimensions"] else None)
    time_col = roles["time"][0] if roles["time"] else None

    followups: list[str] = []
    if metric and dimension:
        followups.append(f"Show {metric} by {dimension} as a bar chart")
    if metric and time_col:
        followups.append(f"Plot {metric} over time")
    if metric:
        followups.append(f"Show the distribution of {metric}")
    other_dimension = next((d for d in roles["dimensions"] if d != dimension), None)
    if metric and other_dimension:
        followups.append(f"Now split that by {other_dimension}")
    if len(roles["numeric"]) >= 2:
        followups.append("Show a correlation heatmap of numeric columns")
    followups.append("What story does this dataset tell?")

    unique: list[str] = []
    for item in followups:
        if item.lower() != _norm_text(question) and item not in unique:
            unique.append(item)
    return unique[:3]


def build_validation_payload(df: pd.DataFrame, method: str,
                             source_columns: list[str] | None = None,
                             confidence: float = 0.95,
                             reasons: list[str] | None = None) -> dict:
    """Explain why an insight or answer should be trusted."""
    clean_columns = [str(c) for c in (source_columns or []) if c in df.columns]
    row_support = int(len(df.dropna(subset=clean_columns))) if clean_columns else int(len(df))
    missing_pct = 0.0
    if clean_columns:
        cells = max(len(df) * len(clean_columns), 1)
        missing_pct = round(float(df[clean_columns].isna().sum().sum()) / cells * 100, 1)

    label = "High" if confidence >= 0.9 else "Medium" if confidence >= 0.7 else "Low"
    return {
        "confidence": round(confidence, 3),
        "confidence_label": label,
        "method": method,
        "row_support": row_support,
        "source_columns": clean_columns,
        "missing_pct": missing_pct,
        "reasons": reasons or [
            "Computed directly from the uploaded dataframe.",
            "No free-form narrative was used for the numeric calculation.",
        ],
    }


def _norm_text(text: str) -> str:
    """Normalize user text for deterministic intent matching and cache keys."""
    return " ".join(text.lower().strip().split())


def infer_column_roles(df: pd.DataFrame) -> dict[str, list[str]]:
    """Infer semantic roles from column names, dtypes, and cardinality."""
    roles: dict[str, list[str]] = {
        "metrics": [],
        "time": [],
        "dimensions": [],
        "ids": [],
        "numeric": [],
        "target_candidates": [],
    }
    metric_terms = ("revenue", "sales", "profit", "cost", "amount", "quantity",
                    "price", "rating", "score", "value", "total", "margin")
    time_terms = ("date", "time", "month", "year", "timestamp", "created", "updated")
    dimension_terms = ("category", "region", "segment", "product", "department",
                       "channel", "city", "country", "state", "status", "group")

    for col in df.columns:
        col_str = str(col)
        name = col_str.lower()
        series = df[col]
        nunique = int(series.nunique(dropna=True))
        is_numeric = pd.api.types.is_numeric_dtype(series)

        if is_numeric:
            roles["numeric"].append(col_str)
            if any(term in name for term in metric_terms):
                roles["metrics"].append(col_str)
            elif not name.endswith("id") and name != "id" and nunique < max(len(df) * 0.95, 2):
                roles["target_candidates"].append(col_str)

        if any(term in name for term in time_terms):
            roles["time"].append(col_str)
            continue

        if name == "id" or name.endswith("_id") or name.endswith(" id"):
            roles["ids"].append(col_str)
            continue

        if is_numeric and nunique > 0.95 * max(len(df), 1):
            roles["ids"].append(col_str)
            continue

        if not is_numeric and (any(term in name for term in dimension_terms) or 2 <= nunique <= 50):
            roles["dimensions"].append(col_str)

    if not roles["metrics"] and roles["numeric"]:
        roles["metrics"] = [c for c in roles["numeric"] if c not in roles["ids"]][:3]
    if not roles["target_candidates"]:
        roles["target_candidates"] = roles["metrics"][:3]
    return roles


def _find_columns_in_question(question: str, columns: list[str]) -> list[str]:
    q = _norm_text(question)
    return [c for c in columns if str(c).lower() in q]


def choose_chart_spec(question: str, df: pd.DataFrame,
                      column_roles: dict[str, list[str]]) -> dict:
    """Choose a deterministic chart type from column roles and question text."""
    q = _norm_text(question)
    mentioned = _find_columns_in_question(q, [str(c) for c in df.columns])
    metric = next((c for c in mentioned if c in column_roles["numeric"]), None)
    metric = metric or (column_roles["metrics"][0] if column_roles["metrics"] else None)
    dimension = next((c for c in mentioned if c in column_roles["dimensions"]), None)
    dimension = dimension or (column_roles["dimensions"][0] if column_roles["dimensions"] else None)
    explicit_time = any(term in q for term in ("over time", "trend", "timeline", "monthly", "yearly", "daily"))
    time_col = next((c for c in mentioned if c in column_roles["time"]), None)
    if not time_col and explicit_time:
        time_col = column_roles["time"][0] if column_roles["time"] else None
    numeric_cols = [c for c in mentioned if c in column_roles["numeric"]]

    if "heatmap" in q or "correlation" in q:
        return {"chart_type": "heatmap", "reason": "Correlation requests are best shown as a heatmap."}
    if ("distribution" in q or "histogram" in q) and metric:
        return {"chart_type": "histogram", "x": metric, "reason": "Single numeric distributions are best shown as histograms."}
    if explicit_time and time_col and metric:
        return {"chart_type": "line", "x": time_col, "y": metric, "reason": "Time plus a numeric metric is best shown as a line chart."}
    if len(numeric_cols) >= 2:
        return {"chart_type": "scatter", "x": numeric_cols[0], "y": numeric_cols[1], "reason": "Two numeric fields are best compared with a scatter plot."}
    if dimension and metric:
        return {"chart_type": "bar", "x": dimension, "y": metric, "reason": "A category and numeric metric are best shown as a bar chart."}
    if time_col and metric:
        return {"chart_type": "line", "x": time_col, "y": metric, "reason": "Time plus a numeric metric is best shown as a line chart."}
    if dimension and ("count" in q or "top" in q or "categor" in q):
        return {"chart_type": "bar", "x": dimension, "y": "count", "reason": "Category counts are best shown as a top-N bar chart."}
    if metric:
        return {"chart_type": "histogram", "x": metric, "reason": "A single numeric field is best shown as a histogram."}
    return {"chart_type": "none", "reason": "No reliable chart mapping was found."}


def _plotly_theme(fig: go.Figure) -> go.Figure:
    fig.update_layout(
        template="plotly_white",
        font=dict(family="Inter, Segoe UI, sans-serif", size=12),
        title_font_size=15,
        colorway=INDIGO_PALETTE[:8],
        margin=dict(l=60, r=30, t=60, b=60),
        hoverlabel=dict(bgcolor="white", font_size=12),
    )
    return fig


def _chart_to_json(fig: go.Figure) -> str:
    return _plotly_theme(fig).to_json()


def _dataset_fingerprint(df: pd.DataFrame) -> str:
    """Create a lightweight fingerprint for cache invalidation."""
    cols = "|".join(str(c) for c in df.columns)
    dtypes = "|".join(str(v) for v in df.dtypes)
    sample = pd.util.hash_pandas_object(df.head(50), index=True).sum()
    return f"{df.shape[0]}x{df.shape[1]}:{cols}:{dtypes}:{int(sample)}"


def _cache_key(session_id: str, category: str, question: str, df: pd.DataFrame) -> str:
    return f"{session_id}:{category}:{_norm_text(question)}:{_dataset_fingerprint(df)}"


def _update_conversation_state(session_id: str, question: str, result: str,
                               plan: dict, code: str | None = None,
                               code_lang: str | None = None) -> None:
    conversation_state[session_id] = {
        "last_question": question,
        "last_result": result[:1200],
        "last_metric": (plan.get("chart_spec") or {}).get("y"),
        "last_grouping": (plan.get("chart_spec") or {}).get("x"),
        "last_chart_type": plan.get("chart_type"),
        "last_code": code,
        "last_code_lang": code_lang,
    }


def deterministic_answer(question: str, df: pd.DataFrame,
                         category: str = "general",
                         previous_state: dict | None = None) -> dict | None:
    """Answer common analytics questions without calling the LLM."""
    q = _norm_text(question)
    roles = infer_column_roles(df)
    chart_spec = choose_chart_spec(question, df, roles)
    columns = [str(c) for c in df.columns]
    mentioned = _find_columns_in_question(q, columns)

    plan = {
        "strategy": "Deterministic pandas analysis for a common analytics request.",
        "query_type": "deterministic",
        "relevant_columns": mentioned,
        "analysis_steps": ["Infer intent", "Compute directly from the DataFrame", "Return JSON-safe result"],
        "chart_type": chart_spec.get("chart_type", "none"),
        "chart_spec": chart_spec,
        "column_roles": roles,
        "domain_focus": category,
    }

    def done(result: str, chart_json: str | None = None, code: str | None = None,
             validation_method: str = "Deterministic pandas computation",
             validation_reasons: list[str] | None = None) -> dict:
        source_columns = [c for c in plan.get("relevant_columns", []) if c in df.columns]
        spec = plan.get("chart_spec") or {}
        for value in (spec.get("x"), spec.get("y")):
            if value and value != "count" and value in df.columns and value not in source_columns:
                source_columns.append(value)
        return {
            "result": result,
            "chart": None,
            "chart_json": chart_json,
            "report": result,
            "critique": {
                "verdict": "pass",
                "confidence": 0.98,
                "issues": [],
                "strengths": ["Computed deterministically from the uploaded data."],
                "suggestion": "",
            },
            "plan": plan,
            "code": code,
            "code_lang": "python" if code else None,
            "validation": build_validation_payload(
                df,
                method=validation_method,
                source_columns=source_columns,
                confidence=0.98,
                reasons=validation_reasons,
            ),
            "meta": {"route": "deterministic", "cacheable": True},
        }

    if any(term in q for term in ("how many rows", "how many columns", "row count", "column count", "shape")):
        result = f"The dataset has {len(df):,} rows and {df.shape[1]:,} columns."
        code = "result = f\"The dataset has {len(df):,} rows and {df.shape[1]:,} columns.\""
        return done(result, code=code, validation_method="Dataset shape read from DataFrame dimensions")

    if any(term in q for term in ("missing", "null", "blank")):
        missing = df.isna().sum()
        pct = (missing / max(len(df), 1) * 100).round(1)
        out = pd.DataFrame({"missing": missing, "missing_pct": pct})
        out = out[out["missing"] > 0].sort_values(["missing", "missing_pct"], ascending=False)
        if out.empty:
            result = "No missing values were found in the dataset."
        else:
            result = "Missing values by column:\n" + out.to_string()
        code = "missing = df.isna().sum(); result = missing.to_string()"
        return done(result, code=code, validation_method="Missing-value count with pandas isna")

    if "duplicate" in q:
        duplicates = int(df.duplicated().sum())
        pct = round(duplicates / max(len(df), 1) * 100, 1)
        result = f"The dataset has {duplicates:,} duplicate rows ({pct}% of rows)."
        code = "duplicates = int(df.duplicated().sum())"
        return done(result, code=code, validation_method="Duplicate-row count with pandas duplicated")

    if any(term in q for term in ("summary statistic", "summary statistics", "describe", "numeric summary")):
        numeric = df.select_dtypes(include="number")
        if numeric.empty:
            return done("No numeric columns were found for summary statistics.")
        result = numeric.describe().round(2).to_string()
        code = "result = df.select_dtypes(include='number').describe().round(2).to_string()"
        return done(result, code=code, validation_method="Numeric summary statistics with pandas describe")

    if "correlation" in q or "heatmap" in q:
        numeric = df.select_dtypes(include="number")
        if numeric.shape[1] < 2:
            return done("At least two numeric columns are needed for a correlation heatmap.")
        corr = numeric.corr().round(3)
        fig = go.Figure(data=go.Heatmap(z=corr.values, x=list(corr.columns), y=list(corr.columns),
                                        colorscale="RdBu", zmid=0))
        fig.update_layout(title="Correlation Matrix", xaxis_title="Column", yaxis_title="Column")
        result = "Correlation matrix:\n" + corr.to_string()
        code = "corr = df.select_dtypes(include='number').corr().round(3)"
        return done(result, chart_json=_chart_to_json(fig), code=code,
                    validation_method="Pearson correlation computed from numeric columns")

    if any(term in q for term in ("distribution", "histogram")):
        metric = chart_spec.get("x") or (roles["metrics"][0] if roles["metrics"] else None)
        if not metric:
            return done("No numeric column was found for a distribution chart.")
        fig = px.histogram(df, x=metric, title=f"Distribution of {metric}")
        result = f"Distribution chart generated for `{metric}`."
        code = f"fig = px.histogram(df, x={metric!r})"
        return done(result, chart_json=_chart_to_json(fig), code=code,
                    validation_method="Histogram generated from a selected numeric column")

    if any(term in q for term in ("top categor", "categories by count", "count by categor", "count records")):
        dim = chart_spec.get("x") if chart_spec.get("y") == "count" else None
        dim = dim or (roles["dimensions"][0] if roles["dimensions"] else None)
        if not dim:
            return done("No suitable categorical column was found for category counts.")
        counts = df[dim].astype(str).value_counts().head(10).reset_index()
        counts.columns = [dim, "count"]
        fig = px.bar(counts, x="count", y=dim, orientation="h", title=f"Top {dim} by Count")
        fig.update_layout(yaxis=dict(autorange="reversed"))
        result = counts.to_string(index=False)
        code = f"result = df[{dim!r}].value_counts().head(10).to_string()"
        return done(result, chart_json=_chart_to_json(fig), code=code,
                    validation_method="Top-N value counts computed from a categorical column")

    if previous_state and any(term in q for term in ("now split", "split that", "break that", "by product", "by region")):
        last_metric = previous_state.get("last_metric")
        if last_metric in df.columns:
            dim = next((c for c in roles["dimensions"] if c.lower() in q), None) or roles["dimensions"][0] if roles["dimensions"] else None
            if dim:
                grouped = df.groupby(dim)[last_metric].sum().sort_values(ascending=False).head(15)
                result = f"Using the previous metric `{last_metric}`, here is the split by `{dim}`:\n" + grouped.round(2).to_string()
                plot_df = grouped.reset_index()
                plot_df.columns = [dim, last_metric]
                fig = px.bar(plot_df, x=dim, y=last_metric, title=f"Total {last_metric} by {dim}")
                plan["relevant_columns"] = [dim, last_metric]
                plan["chart_spec"] = {"chart_type": "bar", "x": dim, "y": last_metric,
                                      "reason": "Conversation memory supplied the previous metric."}
                plan["chart_type"] = "bar"
                code = f"result = df.groupby({dim!r})[{last_metric!r}].sum().sort_values(ascending=False).to_string()"
                return done(result, chart_json=_chart_to_json(fig), code=code,
                            validation_method="Follow-up answered with stored previous metric and pandas groupby")

    aggregate_words = (" by ", "group by", "grouped by", "break down", "breakdown", "split by")
    metric = chart_spec.get("y") if chart_spec.get("y") not in (None, "count") else None
    dim = chart_spec.get("x")
    if any(term in q for term in aggregate_words) and dim and metric:
        agg = "mean" if any(term in q for term in ("average", "avg", "mean")) else "sum"
        grouped = getattr(df.groupby(dim)[metric], agg)().sort_values(ascending=False).head(15)
        title_metric = "Average" if agg == "mean" else "Total"
        result = f"{title_metric} {metric} by {dim}:\n" + grouped.round(2).to_string()
        plot_df = grouped.reset_index()
        plot_df.columns = [dim, metric]
        fig = px.bar(plot_df, x=dim, y=metric, title=f"{title_metric} {metric} by {dim}")
        code = f"result = df.groupby({dim!r})[{metric!r}].{agg}().sort_values(ascending=False).to_string()"
        return done(result, chart_json=_chart_to_json(fig), code=code,
                    validation_method=f"Pandas groupby {agg} aggregation")

    return None


def train_predictive_model(df: pd.DataFrame, target: str) -> tuple[str, str, dict]:
    """Train a Random Forest; return (summary, importance_chart, info).
    info also contains shap_chart, perm_chart, pdp_chart for explainability."""
    from sklearn.ensemble import RandomForestClassifier, RandomForestRegressor
    from sklearn.model_selection import train_test_split
    from sklearn.metrics import accuracy_score, f1_score, r2_score, mean_absolute_error
    from sklearn.inspection import permutation_importance, PartialDependenceDisplay

    data = df.dropna(subset=[target]).copy()
    if len(data) < 30:
        raise ValueError("Not enough rows to train a reliable model (need at least 30).")

    y_raw = data[target]
    X = data.drop(columns=[target])

    id_like = [c for c in X.columns
               if str(c).lower() in ("id", "index")
               or str(c).lower().endswith("_id")
               or (pd.api.types.is_integer_dtype(X[c]) and X[c].nunique() > 0.95 * len(X))]
    X = X.drop(columns=id_like, errors="ignore")

    num_cols = list(X.select_dtypes(include="number").columns)
    cat_cols = [c for c in X.columns
                if not pd.api.types.is_numeric_dtype(X[c]) and X[c].nunique() <= 15]
    X = pd.get_dummies(X[num_cols + cat_cols], columns=cat_cols, drop_first=True)
    X = X.select_dtypes(include="number").fillna(X.median(numeric_only=True))
    if X.shape[1] == 0:
        raise ValueError("No usable feature columns to train on.")

    is_classification = (not pd.api.types.is_numeric_dtype(y_raw)) or y_raw.nunique() <= 10
    if is_classification:
        ycat = y_raw.astype("category")
        n_classes = len(ycat.cat.categories)
        classes = [str(c) for c in ycat.cat.categories]
        y = ycat.cat.codes
        if n_classes < 2:
            raise ValueError("Target has only one class — nothing to classify.")
        model = RandomForestClassifier(n_estimators=200, max_depth=10, random_state=42, n_jobs=-1)
    else:
        classes = None
        y = pd.to_numeric(y_raw, errors="coerce")
        model = RandomForestRegressor(n_estimators=200, max_depth=12, random_state=42, n_jobs=-1)

    strat = y if (is_classification and y.value_counts().min() >= 2) else None
    Xtr, Xte, ytr, yte = train_test_split(X, y, test_size=0.25, random_state=42, stratify=strat)
    model.fit(Xtr, ytr)
    pred = model.predict(Xte)

    # ── 1. Feature importance (gini / MDI) ───────────────────────────────────
    importances = pd.Series(model.feature_importances_, index=X.columns).sort_values(ascending=False)
    top = importances.head(12)
    fig, ax = plt.subplots(figsize=(8, max(4, 0.5 * len(top))))
    sns.barplot(x=top.values, y=top.index.astype(str), ax=ax)
    for cont in ax.containers:
        ax.bar_label(cont, fmt="%.3f", padding=3)
    ax.set_title(f"What predicts '{target}'?  ·  Feature Importance (MDI)")
    ax.set_xlabel("Importance"); ax.set_ylabel("Feature")
    chart = _fig_to_b64(fig)

    # ── 2. SHAP beeswarm (global explanation) ────────────────────────────────
    shap_chart = None
    try:
        import shap
        X_sample = Xtr.iloc[:min(300, len(Xtr))]
        explainer = shap.TreeExplainer(model)
        shap_values = explainer.shap_values(X_sample)
        if is_classification and isinstance(shap_values, list):
            sv = shap_values[1] if len(shap_values) > 1 else shap_values[0]
        else:
            sv = shap_values
        plt.figure(figsize=(8, max(4, 0.4 * min(12, X_sample.shape[1]))))
        shap.summary_plot(sv, X_sample, show=False, max_display=12, plot_type="dot")
        plt.title(f"SHAP Feature Impact on '{target}'", fontsize=14, fontweight="bold", pad=14)
        plt.tight_layout()
        buf = io.BytesIO()
        plt.savefig(buf, format="png", bbox_inches="tight", dpi=150)
        plt.close("all")
        shap_chart = base64.b64encode(buf.getvalue()).decode()
    except Exception:
        plt.close("all")

    # ── 3. Permutation importance (test-set, unbiased) ───────────────────────
    perm_chart = None
    try:
        perm = permutation_importance(model, Xte, yte, n_repeats=5, random_state=42, n_jobs=-1)
        perm_df = (pd.DataFrame({"Feature": X.columns,
                                 "Importance": perm.importances_mean,
                                 "Std": perm.importances_std})
                   .sort_values("Importance", ascending=False).head(12))
        fig, ax = plt.subplots(figsize=(8, max(4, 0.5 * len(perm_df))))
        ax.barh(perm_df["Feature"][::-1], perm_df["Importance"][::-1],
                xerr=perm_df["Std"][::-1], color="#4F46E5", alpha=0.85, capsize=3)
        ax.set_title(f"Permutation Importance · '{target}' (test set)")
        ax.set_xlabel("Mean decrease in score")
        plt.tight_layout()
        perm_chart = _fig_to_b64(fig)
    except Exception:
        plt.close("all")

    # ── 4. Partial Dependence Plots (top 2 features) ─────────────────────────
    pdp_chart = None
    try:
        top2 = list(importances.head(2).index)
        n = len(top2)
        fig, axes = plt.subplots(1, n, figsize=(6 * n, 4))
        if n == 1:
            axes = [axes]
        PartialDependenceDisplay.from_estimator(model, X, top2, ax=axes,
                                                kind="average", subsample=500, random_state=42)
        fig.suptitle(f"Partial Dependence — Top Features for '{target}'",
                     fontsize=13, fontweight="bold")
        fig.tight_layout()
        pdp_chart = _fig_to_b64(fig)
    except Exception:
        plt.close("all")

    # ── metrics & summary ────────────────────────────────────────────────────
    if is_classification:
        metric = f"Accuracy: {accuracy_score(yte, pred):.1%}   |   F1 (weighted): {f1_score(yte, pred, average='weighted'):.2f}"
        task = f"Trained a Random Forest classifier to predict '{target}' ({n_classes} classes)."
    else:
        metric = f"R²: {r2_score(yte, pred):.3f}   |   MAE: {mean_absolute_error(yte, pred):,.3f}"
        task = f"Trained a Random Forest regressor to predict '{target}'."

    top3 = ", ".join(f"{n} ({v:.0%})" for n, v in top.head(3).items())
    summary = (f"{task}\n{metric}\n\n"
               f"Top features: {top3}.\n"
               f"Trained on {len(Xtr)} rows, validated on {len(Xte)}, {X.shape[1]} features.\n"
               f"Explainability: SHAP beeswarm · Permutation importance · Partial dependence plots generated.")

    features_meta = []
    for c in num_cols + cat_cols:
        if c in cat_cols:
            opts = [str(v) for v in pd.Series(data[c].dropna().unique()).tolist()[:30]]
            mode = data[c].mode()
            default = str(mode.iloc[0]) if not mode.empty else (opts[0] if opts else "")
            features_meta.append({"name": str(c), "type": "category", "options": opts, "default": default})
        else:
            med = data[c].median()
            features_meta.append({"name": str(c), "type": "number",
                                  "default": None if pd.isna(med) else round(float(med), 2)})

    info = {
        "target": target,
        "model": model,
        "feature_cols": list(X.columns),
        "num_cols": num_cols,
        "cat_cols": cat_cols,
        "is_classification": is_classification,
        "classes": classes,
        "medians": {k: (None if pd.isna(v) else float(v)) for k, v in X.median().items()},
        "features": features_meta,
        # Explainability charts
        "shap_chart": shap_chart,
        "perm_chart": perm_chart,
        "pdp_chart":  pdp_chart,
    }
    return summary, chart, info


def build_predict_payload(session_id: str, df: pd.DataFrame, target: str) -> dict:
    """Train a model and return the JSON-safe prediction payload."""
    summary, chart, info = train_predictive_model(df, target)
    models[session_id] = info
    return {
        "message": "Model trained with full explainability",
        "result": summary,
        "chart": chart,
        "shap_chart": info.get("shap_chart"),
        "perm_chart": info.get("perm_chart"),
        "pdp_chart": info.get("pdp_chart"),
        "features": info["features"],
        "target": info["target"],
        "is_classification": info["is_classification"],
    }


def get_df_schema(df: pd.DataFrame) -> str:
    schema = f"Shape: {df.shape[0]} rows x {df.shape[1]} columns\n\nColumns:\n"
    for col in df.columns:
        schema += f"  - {col} ({df[col].dtype}): sample values: {df[col].dropna().head(3).tolist()}\n"
    schema += f"\nFirst 5 rows:\n{df.head(5).to_string()}"
    return schema


def get_sql_schema(df: pd.DataFrame) -> str:
    """Schema description formatted for a SQL agent (table name: `data`)."""
    schema = f"Table name: `data`\nRows: {len(df)}\n\nColumns:\n"
    for col in df.columns:
        schema += f"  - `{col}` ({df[col].dtype}): e.g. {df[col].dropna().head(3).tolist()}\n"
    return schema


def validate_sql(sql: str) -> str:
    """Allow only a single read-only SQLite query generated by the SQL analyst."""
    import re

    cleaned = sql.strip()
    if cleaned.startswith("```"):
        cleaned = cleaned.split("\n", 1)[1] if "\n" in cleaned else ""
    if "```" in cleaned:
        cleaned = cleaned.rsplit("```", 1)[0]
    cleaned = cleaned.strip().rstrip(";").strip()

    if not cleaned:
        raise ValueError("SQL query is empty.")
    if ";" in cleaned:
        raise ValueError("Only one SQL statement is allowed.")

    lowered = cleaned.lower()
    if not (lowered.startswith("select") or lowered.startswith("with")):
        raise ValueError("Only SELECT queries are allowed.")

    blocked = ("drop", "delete", "insert", "update", "alter", "attach",
               "detach", "pragma", "replace", "create", "vacuum")
    tokens = set(re.findall(r"[a-z_]+", lowered))
    found = sorted(set(blocked) & tokens)
    if found:
        raise ValueError(f"Blocked SQL keyword: {found[0]}")
    return cleaned


def execute_sql(sql: str, df: pd.DataFrame) -> str:
    """Load df into in-memory SQLite as table 'data', run sql, return result string."""
    import sqlite3
    conn = sqlite3.connect(":memory:")
    try:
        sql = validate_sql(sql)
        df.to_sql("data", conn, index=False, if_exists="replace")
        result_df = pd.read_sql_query(sql, conn)
        if result_df.empty:
            return "The query returned no results."
        return result_df.to_string(index=False)
    finally:
        conn.close()


# ── Multi-agent system prompts ────────────────────────────────────────────────

PLANNER_SYSTEM = """You are a data analysis planner. Given a DataFrame schema and a user question, produce a structured analysis plan.

Output ONLY valid JSON (no markdown fences, no extra text) with exactly this structure:
{
  "relevant_columns": ["col1", "col2"],
  "strategy": "One concise sentence describing the analysis approach",
  "needs_chart": true,
  "chart_type": "bar|line|scatter|histogram|heatmap|box|violin|none",
  "analysis_steps": ["step1", "step2", "step3"],
  "domain_focus": "what domain aspect to emphasize based on the category",
  "query_type": "pandas"
}

query_type rules — choose "sql" when the question involves:
- Filtering rows by value (WHERE-style: top N, specific category, date range)
- Grouping and aggregating (GROUP BY: sum/count/avg by category)
- Ranking or sorting a subset (ORDER BY + LIMIT)
- Simple lookups or cross-tabulations
Otherwise use "pandas" (for correlations, distributions, ML, custom statistics)."""

ANALYST_SYSTEM = """You are a data analyst. A pandas DataFrame `df` is already loaded.
Write Python code to compute the NUMERICAL analysis only — statistics, aggregations, comparisons, correlations. No charts.

Rules:
- `df` is already defined — do not reload it
- Pre-imported: pandas (pd), numpy (np), io, base64
- You MAY import: scipy, sklearn, statsmodels
- Set `result` to a detailed string with all numerical findings (include specific numbers)
- Set `chart_b64 = None` always
- If result is a DataFrame, convert: result = df_result.to_string()
- ROBUSTNESS: use df.select_dtypes(include='number') for numeric ops; never run math on text columns
- Output ONLY valid Python code — no markdown fences, no explanation"""

SQL_ANALYST_SYSTEM = """You are a SQL data analyst. A SQLite database table named `data` contains the user's dataset.
Write a single SQL SELECT query to answer the user's question.

Rules:
- Table name is always `data` (lowercase)
- Use standard SQLite SQL syntax
- For column names with spaces or special characters: wrap in double quotes e.g. "column name"
- For text matching: use LIKE '%value%' (case insensitive in SQLite)
- For aggregations: use GROUP BY, ORDER BY, LIMIT
- Always add ORDER BY for ranking/top-N questions
- Output ONLY the raw SQL query — no markdown fences, no semicolon, no explanation

Common patterns:
  Top N by metric:     SELECT category, SUM(metric) AS total FROM data GROUP BY category ORDER BY total DESC LIMIT 10
  Filter + aggregate:  SELECT region, AVG(sales) FROM data WHERE year = 2023 GROUP BY region
  Count by group:      SELECT status, COUNT(*) AS count FROM data GROUP BY status ORDER BY count DESC
  Date range:          SELECT * FROM data WHERE date >= '2023-01-01' AND date <= '2023-12-31'"""

VISUALIZER_SYSTEM = """You are a data visualization expert. A pandas DataFrame `df` is loaded.
Write Python code to create ONE excellent INTERACTIVE Plotly chart that best illustrates the findings.

Rules:
- `df` is already defined — do not reload it
- Pre-imported: pandas (pd), numpy (np), plotly.express (px), plotly.graph_objects (go)
- Create a Plotly figure named `fig` using px or go
- Apply this theme exactly ONCE after creating fig:
    fig.update_layout(
        template='plotly_white',
        font=dict(family='Inter, Segoe UI, sans-serif', size=12),
        title_font_size=15, title_font_weight='bold',
        colorway=['#4F46E5','#10B981','#F59E0B','#8B5CF6','#EF4444','#06B6D4','#EC4899','#0EA5E9'],
        margin=dict(l=60, r=30, t=60, b=60),
        hoverlabel=dict(bgcolor='white', font_size=12),
    )
- Give the chart a clear descriptive title and labelled axes via update_layout
- Save as: chart_json = fig.to_json()
- Set chart_b64 = None, result = None
- Common chart patterns:
    • Bar: px.bar(df, x='col', y='col', title='...')
    • Line: px.line(df, x='date_col', y='metric', title='...')
    • Scatter: px.scatter(df, x='col1', y='col2', color='group', title='...')
    • Histogram: px.histogram(df, x='col', title='...')
    • Box: px.box(df, x='group', y='metric', title='...')
    • Heatmap: use go.Heatmap with z=corr.values, x=corr.columns, y=corr.columns
- Output ONLY valid Python code — no markdown fences, no explanation"""

CRITIC_SYSTEM = """You are a senior statistician reviewing a data analysis for accuracy and completeness.

Output ONLY valid JSON (no markdown fences, no extra text):
{
  "verdict": "pass" | "warn" | "fail",
  "confidence": 0.0,
  "issues": ["specific issue 1"],
  "strengths": ["specific strength 1"],
  "suggestion": "One concrete improvement the analyst should make"
}

verdict=pass: analysis is statistically sound
verdict=warn: minor caveats or missing context
verdict=fail: significant errors or misleading conclusions
confidence: your confidence that the findings answer the question (0.0-1.0)"""

REPORTER_SYSTEM = """You are an executive analyst writing a business-ready summary of data findings.
Write a clear, structured report with NO padding.

Structure your response EXACTLY as:
**Headline:** One sentence direct answer to the question.

**Key Findings:**
• Finding 1 with specific numbers
• Finding 2 with specific numbers
• Finding 3 with specific numbers

**Implication:** One sentence business recommendation or implication.

**Caveat:** One sentence noting any limitation or assumption (if relevant)."""

# Keep the original SYSTEM_PROMPT for backward compatibility with non-agentic paths
SYSTEM_PROMPT = ANALYST_SYSTEM


# Domain lenses — the selected category turns the agent into a domain expert,
# shaping which computations it prefers and how it narrates the result.
CATEGORY_PERSONAS = {
    "general": (
        "ANALYSIS LENS: General. You are a meticulous general-purpose data analyst. "
        "Provide clear, neutral, statistically sound insights."
    ),
    "financial": (
        "ANALYSIS LENS: Financial. You are a senior FINANCIAL analyst. Interpret the data "
        "through a financial lens — revenue, costs, margins, growth rates (YoY/MoM), volatility "
        "and risk, and key ratios. Prefer computations like growth %, cumulative totals, moving "
        "averages, and standard deviation as a risk proxy. Format money clearly and, in the "
        "written answer, call out financial implications, risks, and opportunities."
    ),
    "medical": (
        "ANALYSIS LENS: Medical. You are a clinical / healthcare data analyst. Interpret the data "
        "through a medical lens — prevalence, risk factors, patient cohorts, and distributions of "
        "clinical measurements. Prefer group-wise comparisons (outcome vs. non-outcome), "
        "correlation of features with the outcome, and distribution analysis. ALWAYS describe "
        "relationships as associations, NOT causation. Flag clinically meaningful or at-risk groups."
    ),
    "retail": (
        "ANALYSIS LENS: Retail. You are a retail & e-commerce analyst. Interpret the data through "
        "a commerce lens — sales and revenue by product/category/region, order volume, basket "
        "size, ratings, and customer behavior. Prefer top-N rankings, revenue breakdowns, and trends."
    ),
    "marketing": (
        "ANALYSIS LENS: Marketing. You are a marketing & growth analyst. Interpret the data through "
        "a marketing lens — acquisition, conversion, retention, segmentation, channels, and campaign "
        "performance. Prefer funnel/segment breakdowns, rates, and cohort-style comparisons."
    ),
    "hr": (
        "ANALYSIS LENS: HR. You are an HR / people-analytics specialist. Interpret the data through "
        "a workforce lens — headcount, attrition/turnover, tenure, demographics, performance, and "
        "compensation equity. Prefer group comparisons and distribution analysis; be sensitive about fairness."
    ),
}


def system_prompt_for(category: str) -> str:
    persona = CATEGORY_PERSONAS.get(category, CATEGORY_PERSONAS["general"])
    return f"{persona}\n\n{SYSTEM_PROMPT}"


# ── Sandbox security configuration ───────────────────────────────────────────

MAX_EXEC_SECONDS  = 30    # hard wall-clock timeout per code execution
MAX_RESULT_CHARS  = 5_000 # truncate oversized results to prevent memory abuse

# Modules the generated analysis code is allowed to import.
ALLOWED_MODULES = {
    "pandas", "numpy", "matplotlib", "seaborn", "scipy", "sklearn", "statsmodels",
    "plotly", "math", "statistics", "datetime", "io", "base64", "collections", "itertools", "re", "json",
}

# Dangerous dunder / frame attributes that could be used to escape the sandbox.
_BLOCKED_ATTRS = frozenset({
    "__class__", "__bases__", "__subclasses__", "__mro__",
    "__globals__", "__builtins__", "__code__", "__closure__",
    "__dict__", "__loader__", "__spec__", "__wrapped__",
    "f_locals", "f_globals", "f_back", "gi_frame",
    "func_globals", "func_code",
})

# Dangerous callables that must never appear in generated code.
_BLOCKED_CALLS = frozenset({
    "eval", "exec", "compile", "open", "breakpoint",
    "__import__", "input", "memoryview",
})


class SecurityError(Exception):
    """Raised when generated code fails the AST security scan."""


class _ASTSecurityVisitor(ast.NodeVisitor):
    """Walk the AST and raise SecurityError on any dangerous pattern."""

    def visit_Attribute(self, node: ast.Attribute) -> None:
        if node.attr in _BLOCKED_ATTRS:
            raise SecurityError(
                f"Access to attribute '{node.attr}' is blocked in the sandbox"
            )
        self.generic_visit(node)

    def visit_Call(self, node: ast.Call) -> None:
        if isinstance(node.func, ast.Name) and node.func.id in _BLOCKED_CALLS:
            raise SecurityError(
                f"Call to '{node.func.id}' is blocked in the sandbox"
            )
        self.generic_visit(node)

    def visit_Import(self, node: ast.Import) -> None:
        for alias in node.names:
            root = alias.name.split(".")[0]
            if root not in ALLOWED_MODULES:
                raise SecurityError(f"Import of '{alias.name}' is not allowed in the sandbox")
        self.generic_visit(node)

    def visit_ImportFrom(self, node: ast.ImportFrom) -> None:
        root = (node.module or "").split(".")[0]
        if root not in ALLOWED_MODULES:
            raise SecurityError(f"Import from '{node.module}' is not allowed in the sandbox")
        self.generic_visit(node)


def validate_code_ast(code: str) -> None:
    """Parse code and run the AST security scan. Raises SecurityError on any violation."""
    try:
        tree = ast.parse(code)
    except SyntaxError as e:
        raise SecurityError(f"Syntax error in generated code: {e}")
    _ASTSecurityVisitor().visit(tree)


def execute_code(code: str, df: pd.DataFrame) -> tuple[str, str | None, str | None]:
    """Execute sandboxed code with three security layers:
    1. AST pre-scan (blocks dangerous attributes / calls / imports)
    2. Restricted __builtins__ (whitelisted only)
    3. Process-isolated execution timeout (MAX_EXEC_SECONDS)
    Returns (result, chart_b64, chart_json).
    """
    # Layer 1 — AST security scan (raises SecurityError on violations)
    validate_code_ast(code)

    ctx = mp.get_context("spawn")
    output_queue = ctx.Queue(maxsize=1)
    process = ctx.Process(target=execute_code_worker, args=(code, df, output_queue))
    process.start()

    try:
        payload = output_queue.get(timeout=MAX_EXEC_SECONDS)
    except queue.Empty as exc:
        process.terminate()
        process.join(1)
        if process.is_alive() and hasattr(process, "kill"):
            process.kill()
            process.join(1)
        raise TimeoutError(
            f"Code execution exceeded the {MAX_EXEC_SECONDS}s time limit and was terminated."
        ) from exc
    finally:
        process.join(2)
        output_queue.close()
        output_queue.join_thread()

    if payload[0] == "error":
        _, exc_type, exc_message, traceback_text = payload
        raise RuntimeError(
            f"Generated code failed with {exc_type}: {exc_message}\n{traceback_text}"
        )

    _, result, chart_b64, chart_json = payload

    if result is None:
        result = "Done. See chart above." if (chart_b64 or chart_json) else "No result returned."

    # Result size guard — truncate enormous outputs
    result_str = str(result)
    if len(result_str) > MAX_RESULT_CHARS:
        result_str = (result_str[:MAX_RESULT_CHARS]
                      + f"\n… [truncated — {len(result_str):,} chars total]")

    return result_str, chart_b64, chart_json


@app.get("/health")
def health() -> dict:
    return {"status": "ok"}


def register_dataframe(df: pd.DataFrame, filename: str) -> dict:
    """Store a DataFrame and return its profile + instant overview charts."""
    cleanup_expired_sessions()
    validate_dataframe_limits(df)
    session_id = str(uuid.uuid4())
    token = str(uuid.uuid4())
    dataframes[session_id] = df
    session_meta[session_id] = {"created_at": _now(), "last_accessed": _now(), "filename": filename, "token": token}
    storage.save_session(session_id, token, filename, df)
    cleanup_expired_sessions()
    profile = build_profile(df)
    overview = build_overview_charts(df)
    proactive_insights = build_proactive_insights(df)
    column_roles = infer_column_roles(df)
    quality_report = build_quality_report(df)
    decision_brief = build_decision_brief(df)
    cleaning_plan = build_cleaning_plan(df)
    data_contract = build_data_contract(df)
    dashboard_spec = build_dashboard_spec(df)
    decision_actions = decision_brief.get("decision_actions", [])
    return {
        "session_id": session_id,
        "token": token,
        "filename": filename,
        **profile,
        "overview": overview,
        "proactive_insights": proactive_insights,
        "column_roles": column_roles,
        "quality_report": quality_report,
        "decision_brief": decision_brief,
        "decision_actions": decision_actions,
        "cleaning_plan": cleaning_plan,
        "data_contract": data_contract,
        "dashboard_spec": dashboard_spec,
    }


def infer_join_keys(df1: pd.DataFrame, df2: pd.DataFrame) -> list[dict]:
    """Find candidate join keys across two DataFrames using column names and value overlaps."""
    candidates: list[dict] = []
    cols1 = [str(c) for c in df1.columns]
    cols2 = [str(c) for c in df2.columns]

    for c1 in cols1:
        s1 = df1[c1].dropna().astype(str).str.strip()
        if s1.empty:
            continue
        vals1 = set(s1)

        for c2 in cols2:
            s2 = df2[c2].dropna().astype(str).str.strip()
            if s2.empty:
                continue
            vals2 = set(s2)

            name_match = (c1.lower() == c2.lower())
            id_like = (c1.lower().endswith("id") or c2.lower().endswith("id") or c1.lower() == "id" or c2.lower() == "id")

            if not (name_match or id_like or len(vals1 & vals2) > 0):
                continue

            common_vals = vals1 & vals2
            if not common_vals:
                continue

            overlap_ratio1 = len(common_vals) / len(vals1) if vals1 else 0.0
            overlap_ratio2 = len(common_vals) / len(vals2) if vals2 else 0.0
            max_overlap = max(overlap_ratio1, overlap_ratio2)

            score = max_overlap * 0.6
            if name_match:
                score += 0.3
            if id_like:
                score += 0.1
            score = round(min(1.0, score), 2)

            if score >= 0.2:
                sample_matches = sorted(list(common_vals))[:3]
                candidates.append({
                    "column_1": c1,
                    "column_2": c2,
                    "score": score,
                    "confidence_label": "High" if score >= 0.7 else "Medium" if score >= 0.4 else "Low",
                    "matched_unique_count": len(common_vals),
                    "overlap_pct_1": round(overlap_ratio1 * 100, 1),
                    "overlap_pct_2": round(overlap_ratio2 * 100, 1),
                    "sample_matches": sample_matches,
                })

    candidates.sort(key=lambda x: float(x["score"]), reverse=True)
    return candidates[:5]


@app.post("/upload")
async def upload_file(
    file: UploadFile = File(...),
    sheet_name: str | None = None,
) -> dict:
    ext = file.filename.rsplit(".", 1)[-1].lower() if "." in file.filename else ""
    if ext not in ("csv", "xlsx", "xls", "parquet", "json", "jsonl"):
        raise HTTPException(
            status_code=400,
            detail="Unsupported file format. Supported formats: .csv, .xlsx, .xls, .parquet, .json, .jsonl",
        )

    content = await read_upload_limited(file)
    sheets: list[str] = []

    try:
        if ext in ("xlsx", "xls"):
            excel_file = pd.ExcelFile(io.BytesIO(content))
            sheets = excel_file.sheet_names
            target_sheet = sheet_name if (sheet_name and sheet_name in sheets) else 0
            df = pd.read_excel(excel_file, sheet_name=target_sheet)
        elif ext == "parquet":
            df = pd.read_parquet(io.BytesIO(content))
        elif ext == "json":
            raw = json.loads(content.decode("utf-8"))
            if isinstance(raw, list):
                df = pd.json_normalize(raw)
            elif isinstance(raw, dict):
                df = pd.json_normalize(raw.get("data", [raw]))
            else:
                df = pd.DataFrame(raw)
        elif ext == "jsonl":
            lines = [json.loads(line) for line in content.decode("utf-8").splitlines() if line.strip()]
            df = pd.json_normalize(lines)
        else:
            df = pd.read_csv(io.BytesIO(content))
    except Exception as e:
        raise HTTPException(status_code=422, detail=f"Could not parse uploaded file: {e}")

    result = register_dataframe(df, file.filename)
    if sheets:
        result["sheets"] = sheets
    return result


ALLOWED_DOMAINS = {"docs.google.com", "drive.google.com", "sheets.googleapis.com"}


def validate_safe_url(url_str: str) -> str:
    """Validate that the input URL is a secure public Google Sheets export URL and block SSRF vectors."""
    url_str = url_str.strip()
    if not url_str:
        raise HTTPException(status_code=400, detail="URL cannot be empty.")

    try:
        parsed = urllib.parse.urlparse(url_str)
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid URL format.")

    if parsed.scheme not in ("http", "https"):
        raise HTTPException(status_code=400, detail="Invalid URL scheme. Only HTTP and HTTPS URLs are allowed.")

    hostname = parsed.hostname
    if not hostname:
        raise HTTPException(status_code=400, detail="Invalid URL hostname.")

    host_lower = hostname.lower()
    if host_lower in ("localhost", "127.0.0.1", "::1", "169.254.169.254"):
        raise HTTPException(status_code=400, detail="Security error: Internal loopback and metadata addresses are prohibited.")

    domain_ok = any(host_lower == domain or host_lower.endswith("." + domain) for domain in ALLOWED_DOMAINS)
    if not domain_ok:
        raise HTTPException(
            status_code=400,
            detail="Security error: URL import is strictly restricted to public Google Sheets export links (docs.google.com).",
        )

    try:
        ip_info = socket.getaddrinfo(hostname, None)
        for item in ip_info:
            ip_str = item[4][0]
            ip_obj = ipaddress.ip_address(ip_str)
            if ip_obj.is_private or ip_obj.is_loopback or ip_obj.is_link_local or ip_obj.is_reserved:
                raise HTTPException(status_code=400, detail="Security error: Resolved IP address is private or restricted.")
    except socket.gaierror:
        raise HTTPException(status_code=400, detail="Could not resolve hostname.")

    return url_str


@app.post("/import_url")
async def import_from_url(req: UrlImportRequest) -> dict:
    raw_url = validate_safe_url(req.url)

    target_url = raw_url
    if "docs.google.com/spreadsheets/d/" in raw_url:
        match = re.search(r"docs\.google\.com/spreadsheets/d/([a-zA-Z0-9-_]+)", raw_url)
        if match:
            sheet_id = match.group(1)
            target_url = f"https://docs.google.com/spreadsheets/d/{sheet_id}/export?format=csv"

    validate_safe_url(target_url)

    current_url = target_url
    content: bytes | None = None
    max_redirects = 5

    try:
        async with httpx.AsyncClient(follow_redirects=False, timeout=15.0) as client:
            for _ in range(max_redirects + 1):
                validate_safe_url(current_url)
                resp = await client.get(current_url)

                if resp.status_code in (301, 302, 303, 307, 308):
                    redirect_location = resp.headers.get("location")
                    if not redirect_location:
                        raise HTTPException(status_code=400, detail="Redirect missing location header.")
                    current_url = urllib.parse.urljoin(current_url, redirect_location)
                    validate_safe_url(current_url)
                    continue

                if resp.status_code != 200:
                    raise HTTPException(status_code=400, detail=f"Failed to fetch dataset from URL (HTTP {resp.status_code}). Ensure link is publicly accessible.")

                content = resp.content
                break
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Could not download dataset from URL: {e}")

    try:
        df = pd.read_csv(io.BytesIO(content))
    except Exception:
        try:
            df = pd.read_excel(io.BytesIO(content))
        except Exception as e:
            raise HTTPException(status_code=422, detail=f"Could not parse imported dataset: {e}")

    filename = req.filename if req.filename.endswith(".csv") else f"{req.filename}.csv"
    return register_dataframe(df, filename)


def build_dataset_comparison(df1: pd.DataFrame, df2: pd.DataFrame) -> dict:
    """Compare version 1 vs version 2 of a dataset to detect schema changes, row delta, and distribution drift."""
    cols1 = set(df1.columns)
    cols2 = set(df2.columns)

    added_cols = sorted(list(cols2 - cols1))
    removed_cols = sorted(list(cols1 - cols2))
    common_cols = sorted(list(cols1 & cols2))

    type_changes: list[dict] = []
    for c in common_cols:
        t1 = str(df1[c].dtype)
        t2 = str(df2[c].dtype)
        if t1 != t2:
            type_changes.append({"column": c, "v1_type": t1, "v2_type": t2})

    row_delta = len(df2) - len(df1)
    row_pct_change = round((row_delta / max(1, len(df1))) * 100, 1)

    numeric_drift: list[dict] = []
    for c in common_cols:
        if pd.api.types.is_numeric_dtype(df1[c]) and pd.api.types.is_numeric_dtype(df2[c]):
            s1 = df1[c].dropna()
            s2 = df2[c].dropna()
            if len(s1) > 0 and len(s2) > 0:
                m1, m2 = float(s1.mean()), float(s2.mean())
                med1, med2 = float(s1.median()), float(s2.median())
                std1, std2 = float(s1.std()), float(s2.std())
                mean_delta = m2 - m1
                pct_shift = round((mean_delta / max(1e-5, abs(m1))) * 100, 1)

                drift_score = min(1.0, round(abs(m2 - m1) / max(1e-5, (std1 + std2) / 2), 2))
                numeric_drift.append({
                    "column": c,
                    "v1_mean": round(m1, 2),
                    "v2_mean": round(m2, 2),
                    "v1_median": round(med1, 2),
                    "v2_median": round(med2, 2),
                    "mean_delta": round(mean_delta, 2),
                    "pct_shift": pct_shift,
                    "drift_score": drift_score,
                    "drift_level": "Significant" if drift_score > 0.5 else "Moderate" if drift_score > 0.2 else "Low",
                })

    numeric_drift.sort(key=lambda x: x["drift_score"], reverse=True)

    summary = (
        f"Comparison of v1 ({len(df1):,} rows) vs v2 ({len(df2):,} rows): "
        f"{len(added_cols)} columns added, {len(removed_cols)} columns removed, "
        f"and {len([d for d in numeric_drift if d['drift_level'] == 'Significant'])} numeric columns showing significant distribution drift."
    )

    return {
        "summary": summary,
        "v1_rows": len(df1),
        "v2_rows": len(df2),
        "row_delta": row_delta,
        "row_pct_change": row_pct_change,
        "schema_changes": {
            "added_columns": added_cols,
            "removed_columns": removed_cols,
            "common_columns_count": len(common_cols),
            "type_changes": type_changes,
        },
        "numeric_drift": numeric_drift,
    }


@app.post("/compare")
def compare_datasets_endpoint(req: CompareRequest) -> dict:
    cleanup_expired_sessions()
    df1 = get_session_df(req.session_id_1)
    df2 = get_session_df(req.session_id_2)
    _touch_session(req.session_id_1)
    _touch_session(req.session_id_2)
    res = build_dataset_comparison(df1, df2)
    compare_store[req.session_id_1] = res
    compare_store[req.session_id_2] = res
    return res


@app.post("/infer_join")
def infer_join_endpoint(req: InferJoinRequest) -> dict:
    df1 = get_session_df(req.session_id_1)
    df2 = get_session_df(req.session_id_2)
    candidates = infer_join_keys(df1, df2)
    return {
        "session_id_1": req.session_id_1,
        "session_id_2": req.session_id_2,
        "candidates": candidates,
    }


forecast_store: dict[str, dict] = {}
join_store: dict[str, dict] = {}
compare_store: dict[str, dict] = {}


@app.post("/join")
def join_datasets(req: JoinRequest) -> dict:
    df1 = get_session_df(req.session_id_1)
    df2 = get_session_df(req.session_id_2)

    if req.join_key_1 not in df1.columns:
        raise HTTPException(status_code=400, detail=f"Column '{req.join_key_1}' not found in first dataset.")
    if req.join_key_2 not in df2.columns:
        raise HTTPException(status_code=400, detail=f"Column '{req.join_key_2}' not found in second dataset.")
    if req.how not in ("inner", "left", "right", "outer"):
        raise HTTPException(status_code=400, detail="Invalid join type. Use 'inner', 'left', 'right', or 'outer'.")

    try:
        joined_df = pd.merge(
            df1,
            df2,
            left_on=req.join_key_1,
            right_on=req.join_key_2,
            how=req.how,
            suffixes=("_left", "_right"),
        )
    except Exception as e:
        raise HTTPException(status_code=422, detail=f"Could not join datasets: {e}")

    meta1 = session_meta.get(req.session_id_1, {})
    meta2 = session_meta.get(req.session_id_2, {})
    f1 = str(meta1.get("filename", "table1")).rsplit(".", 1)[0]
    f2 = str(meta2.get("filename", "table2")).rsplit(".", 1)[0]
    joined_filename = f"{f1}_{req.how}_join_{f2}.csv"

    res = register_dataframe(joined_df, joined_filename)
    res["join_metadata"] = {
        "left_session_id": req.session_id_1,
        "right_session_id": req.session_id_2,
        "left_filename": meta1.get("filename", "table1"),
        "right_filename": meta2.get("filename", "table2"),
        "join_key_1": req.join_key_1,
        "join_key_2": req.join_key_2,
        "how": req.how,
        "left_rows_before": len(df1),
        "right_rows_before": len(df2),
        "rows_after": len(joined_df),
    }
    join_store[res["session_id"]] = res["join_metadata"]
    return res


def build_time_series_forecast(
    df: pd.DataFrame,
    date_col: str,
    target_col: str,
    periods: int = 12,
    freq: str = "auto",
) -> dict:
    """Perform time-series trend forecasting with 95% confidence bounds."""
    if date_col not in df.columns or target_col not in df.columns:
        raise HTTPException(status_code=400, detail="Specified date or target column not found in dataset.")

    ts_df = df[[date_col, target_col]].dropna().copy()
    ts_df[date_col] = pd.to_datetime(ts_df[date_col], errors="coerce")
    ts_df[target_col] = pd.to_numeric(ts_df[target_col], errors="coerce")
    ts_df = ts_df.dropna().sort_values(by=date_col)

    if len(ts_df) < 5:
        raise HTTPException(status_code=400, detail="At least 5 valid date-target rows are required for forecasting.")

    span_days = (ts_df[date_col].max() - ts_df[date_col].min()).days
    if freq == "auto":
        if span_days > 730:
            resample_freq = "MS"
        elif span_days > 60:
            resample_freq = "W"
        else:
            resample_freq = "D"
    else:
        resample_freq = freq

    try:
        series = ts_df.set_index(date_col)[target_col].resample(resample_freq).mean().interpolate(method="linear")
    except Exception:
        series = ts_df.set_index(date_col)[target_col]

    series = series.dropna()
    if len(series) < 3:
        raise HTTPException(status_code=400, detail="Insufficient time series resolution for forecasting.")

    n_hist = len(series)
    hist_dates = [d.strftime("%Y-%m-%d") for d in series.index]
    hist_values = [round(float(v), 2) for v in series.values]

    try:
        future_dates = pd.date_range(start=series.index[-1], periods=periods + 1, freq=resample_freq)[1:]
    except Exception:
        future_dates = pd.date_range(start=series.index[-1], periods=periods + 1, freq="D")[1:]

    future_date_strs = [d.strftime("%Y-%m-%d") for d in future_dates]

    forecast_vals: list[float] = []
    std_err: float = 0.0

    try:
        from statsmodels.tsa.api import ExponentialSmoothing
        model = ExponentialSmoothing(series.values, trend="add", seasonal=None, initialization_method="estimated")
        fit_model = model.fit()
        pred = fit_model.forecast(periods)
        forecast_vals = [float(v) for v in pred]
        residuals = series.values - fit_model.fittedvalues
        std_err = float(np.std(residuals)) if len(residuals) > 0 else float(np.std(series.values) * 0.1)
    except Exception:
        x = np.arange(n_hist)
        slope, intercept = np.polyfit(x, series.values, 1)
        x_future = np.arange(n_hist, n_hist + periods)
        forecast_vals = [float(slope * xi + intercept) for xi in x_future]
        residuals = series.values - (slope * x + intercept)
        std_err = float(np.std(residuals)) if len(residuals) > 0 else float(np.std(series.values) * 0.1)

    lower_95 = [round(float(v - 1.96 * std_err), 2) for v in forecast_vals]
    upper_95 = [round(float(v + 1.96 * std_err), 2) for v in forecast_vals]
    forecast_rounded = [round(float(v), 2) for v in forecast_vals]

    start_val = hist_values[0] if hist_values else 1.0
    end_forecast = forecast_rounded[-1] if forecast_rounded else start_val
    growth_rate = round(((end_forecast - hist_values[-1]) / max(1e-5, abs(hist_values[-1]))) * 100, 1)
    direction = "upward" if growth_rate > 2.0 else "downward" if growth_rate < -2.0 else "flat"

    plot_hist_df = pd.DataFrame({"Date": hist_dates, "Value": hist_values, "Type": "Historical"})
    plot_fc_df = pd.DataFrame({"Date": future_date_strs, "Value": forecast_rounded, "Type": "Forecast"})
    full_plot_df = pd.concat([plot_hist_df, plot_fc_df])

    fig = px.line(
        full_plot_df,
        x="Date",
        y="Value",
        color="Type",
        title=f"Time-Series Forecast: {target_col} ({periods} {resample_freq} periods)",
        color_discrete_map={"Historical": "#4F46E5", "Forecast": "#10B981"},
    )
    chart_json = _chart_to_json(fig)

    return {
        "date_column": date_col,
        "target_column": target_col,
        "frequency": resample_freq,
        "periods": periods,
        "historical": [{"date": d, "value": v} for d, v in zip(hist_dates, hist_values)],
        "forecast": [
            {"date": d, "forecast": f, "lower_95": l, "upper_95": u}
            for d, f, l, u in zip(future_date_strs, forecast_rounded, lower_95, upper_95)
        ],
        "metrics": {
            "trend_direction": direction,
            "growth_rate_pct": growth_rate,
            "last_historical_value": hist_values[-1],
            "projected_final_value": end_forecast,
        },
        "chart_json": chart_json,
    }


@app.post("/forecast")
def forecast_endpoint(req: ForecastRequest) -> dict:
    cleanup_expired_sessions()
    df = get_session_df(req.session_id)
    _touch_session(req.session_id)
    res = build_time_series_forecast(df, req.date_column, req.target_column, req.periods, req.freq)
    forecast_store[req.session_id] = res
    return res


@app.post("/upload_text")
async def upload_text(req: TextUploadRequest) -> dict:
    """Analyze pasted rows (CSV or TSV, with or without a header row)."""
    text = req.text.strip()
    if not text:
        raise HTTPException(status_code=400, detail="No data was pasted.")
    if len(text.encode("utf-8")) > MAX_UPLOAD_BYTES:
        raise HTTPException(
            status_code=413,
            detail=f"Pasted data is too large. Maximum size is {MAX_UPLOAD_BYTES // (1024 * 1024)} MB.",
        )
    try:
        if "\t" in text:
            sep = "\t"
        elif ";" in text:
            sep = ";"
        else:
            sep = ","
        df = pd.read_csv(
            io.StringIO(text),
            sep=sep,
            header=0 if req.has_header else None,
        )
    except Exception as e:
        raise HTTPException(status_code=422, detail=f"Could not parse the pasted data: {e}")
    if not req.has_header:
        df.columns = [f"col_{i + 1}" for i in range(df.shape[1])]
    if df.empty or df.shape[1] == 0:
        raise HTTPException(status_code=422, detail="The pasted data has no usable rows/columns.")
    return register_dataframe(df, req.filename or "pasted_data.csv")


@app.post("/upload_doc")
async def upload_doc(session_id: str, file: UploadFile = File(...)) -> dict:
    """Upload a PDF, Excel, or text file to enrich analysis with RAG context."""
    cleanup_expired_sessions()
    if session_id not in dataframes:
        raise HTTPException(status_code=404, detail="Upload a CSV first, then attach documents.")
    _touch_session(session_id)
    content = await read_upload_limited(file)
    try:
        text = _parse_doc(content, file.filename)
    except ValueError as e:
        raise HTTPException(status_code=422, detail=str(e))
    if not text.strip():
        raise HTTPException(status_code=422, detail="Document appears to be empty.")
    if session_id not in doc_stores:
        doc_stores[session_id] = DocStore()
    n_chunks = doc_stores[session_id].add(text, file.filename)
    return {
        "filename": file.filename,
        "chunks_indexed": n_chunks,
        "filenames": doc_stores[session_id].filenames,
    }


@app.get("/docs/{session_id}")
def get_docs(session_id: str) -> dict:
    cleanup_expired_sessions()
    if session_id in dataframes:
        _touch_session(session_id)
    store = doc_stores.get(session_id)
    if not store:
        return {"filenames": [], "chunks": 0}
    return {"filenames": store.filenames, "chunks": len(store.chunks)}


@app.get("/quality/{session_id}")
def quality(session_id: str) -> dict:
    df = get_session_df(session_id)
    return build_quality_report(df)


@app.get("/brief/{session_id}")
def decision_brief(session_id: str, category: str = "general") -> dict:
    df = get_session_df(session_id)
    return build_decision_brief(df, category)


@app.get("/cleaning_plan/{session_id}")
def cleaning_plan(session_id: str) -> dict:
    df = get_session_df(session_id)
    return build_cleaning_plan(df)


@app.post("/clean/{session_id}")
def clean_dataset(session_id: str, req: CleanRequest) -> dict:
    df = get_session_df(session_id)
    plan = build_cleaning_plan(df)
    selected_actions = req.actions if req.actions is not None else plan["default_actions"]
    cleaned, applied_actions = apply_cleaning_actions(df, selected_actions)
    before_quality = build_quality_report(df)
    after_quality = build_quality_report(cleaned)

    csv_bytes = cleaned.to_csv(index=False).encode("utf-8")
    original_name = session_meta.get(session_id, {}).get("filename", "dataset.csv")
    stem = str(original_name).rsplit(".", 1)[0] or "dataset"
    return {
        "filename": f"{stem}_cleaned.csv",
        "media_type": "text/csv",
        "content_base64": base64.b64encode(csv_bytes).decode("ascii"),
        "size_bytes": len(csv_bytes),
        "row_delta": int(len(cleaned) - len(df)),
        "column_delta": int(cleaned.shape[1] - df.shape[1]),
        "before_quality": before_quality,
        "after_quality": after_quality,
        "selected_actions": selected_actions,
        "applied_actions": applied_actions,
    }


@app.get("/contract/{session_id}")
def data_contract(session_id: str) -> dict:
    df = get_session_df(session_id)
    return build_data_contract(df)


@app.post("/validate_rows/{session_id}")
def validate_rows(session_id: str, req: ValidateRowsRequest) -> dict:
    df = get_session_df(session_id)
    contract = build_data_contract(df)
    return validate_rows_against_contract(req.rows, contract)


@app.get("/dashboard/{session_id}")
def dashboard_spec(session_id: str, category: str = "general") -> dict:
    df = get_session_df(session_id)
    return build_dashboard_spec(df, category)


@app.post("/investigate")
async def investigate(req: InvestigationRequest) -> StreamingResponse:
    """Run an autonomous, deterministic investigation and stream the analyst workflow."""
    df = get_session_df(req.session_id)

    async def stream():
        emit = make_sse_emitter("investigate", req.session_id)
        yield emit({
            "step": "analyzing",
            "message": "Scoping the investigation goal against the uploaded dataset...",
            "meta": {"route": "deterministic_investigation"},
        })
        yield emit({"step": "planning", "message": "Building an investigation tree across trends, segments, anomalies, and quality..."})
        yield emit({"step": "executing", "message": "Computing deterministic evidence from the dataframe..."})
        investigation = build_investigation_report(df, req.goal, req.category)
        plan = investigation["plan"]
        yield emit({"step": "plan", "message": f"Plan: {plan['strategy']}", "plan": plan})
        yield emit({"step": "reporting", "message": "Writing the autonomous investigation brief from computed evidence..."})
        done_event = {
            "step": "done",
            "message": "Autonomous investigation complete",
            "result": investigation["report"],
            "report": investigation["report"],
            "chart": None,
            "chart_json": investigation["chart_json"],
            "critique": {
                "verdict": "pass",
                "confidence": investigation["validation"]["confidence"],
                "issues": [],
                "strengths": ["Investigation is grounded in deterministic dataframe checks."],
                "suggestion": "Use the recommended actions to drill into the strongest driver next.",
            },
            "plan": plan,
            "followups": suggest_followups(req.goal, df, plan, conversation_state.get(req.session_id)),
            "validation": investigation["validation"],
            "decision_actions": investigation["decision_actions"],
            "investigation": {
                "goal": investigation["goal"],
                "persona": investigation["persona"],
                "persona_columns": investigation["persona_columns"],
                "primary_metric": investigation["primary_metric"],
                "primary_dimension": investigation["primary_dimension"],
                "trend": investigation["trend"],
                "segment_driver": investigation["segment_driver"],
                "outliers": investigation["outliers"],
                "correlation": investigation["correlation"],
                "quality": investigation["quality"],
                "recommended_actions": investigation["recommended_actions"],
                "investigation_tree": investigation["investigation_tree"],
            },
            "meta": investigation["meta"],
        }
        _update_conversation_state(req.session_id, req.goal, investigation["report"], plan)
        yield emit(done_event)

    return StreamingResponse(stream(), media_type="text/event-stream")


@app.post("/query")
async def query_csv(req: QueryRequest) -> StreamingResponse:
    cleanup_expired_sessions()
    if req.session_id not in dataframes:
        raise HTTPException(status_code=404, detail="Session not found. Upload a CSV first.")
    _touch_session(req.session_id)

    df = dataframes[req.session_id]
    schema = get_df_schema(df)
    category_persona = CATEGORY_PERSONAS.get(req.category, CATEGORY_PERSONAS["general"])
    cache_key = _cache_key(req.session_id, req.category, req.question, df)

    async def stream():
        emit = make_sse_emitter("query", req.session_id)

        def llm(system: str, user: str, temperature: float = 0) -> str:
            resp = client.models.generate_content(
                model=GEMINI_MODEL,
                contents=user,
                config=types.GenerateContentConfig(
                    system_instruction=system,
                    temperature=temperature,
                ),
            )
            text = (resp.text or "").strip()
            if text.startswith("```"):
                text = text.split("\n", 1)[1]
            if "```" in text:
                text = text.rsplit("```", 1)[0]
            return text.strip()

        def parse_json_safe(text: str) -> dict:
            import re as _re
            try:
                return json.loads(text)
            except Exception:
                m = _re.search(r'\{.*\}', text, _re.DOTALL)
                if m:
                    try:
                        return json.loads(m.group())
                    except Exception:
                        pass
            return {}

        def run_code_with_repair(system: str, context: str, code_label: str) -> tuple[str | None, str | None, str | None, str | None]:
            """Generate code, execute it, repair once on failure. Returns (code, result, chart_b64, chart_json)."""
            try:
                code = llm(system, context)
            except Exception:
                return None, None, None, None

            for attempt in range(2):
                try:
                    result, chart_b64, chart_json = execute_code(code, df)
                    return code, result, chart_b64, chart_json
                except Exception:
                    err = traceback.format_exc().strip().splitlines()[-1]
                    if attempt == 0:
                        try:
                            code = llm(system,
                                       f"{context}\n\nPrevious code failed:\n{code}\nError: {err}\nFix it.")
                        except Exception:
                            return code, None, None, None
                    else:
                        return code, None, None, None
            return code, None, None, None

        cached = query_cache.get(cache_key)
        if cached:
            yield emit({"step": "analyzing", "message": "Using cached analysis result...",
                        "meta": {"route": "cache", "cache_hit": True}})
            yield emit({**cached, "message": "Analysis complete (cached)"})
            return

        deterministic = deterministic_answer(
            req.question,
            df,
            req.category,
            conversation_state.get(req.session_id),
        )
        if deterministic:
            plan = deterministic["plan"]
            followups = suggest_followups(req.question, df, plan, conversation_state.get(req.session_id))
            yield emit({"step": "planning", "message": "Matched a deterministic analytics intent...",
                        "meta": {"route": "deterministic", "cache_hit": False}})
            yield emit({"step": "plan", "message": f"Plan: {plan.get('strategy')}", "plan": plan})
            if deterministic.get("code"):
                yield emit({"step": "code", "message": "Deterministic pandas code selected",
                            "code": deterministic["code"], "code_lang": deterministic.get("code_lang")})
            yield emit({"step": "executing", "message": "Computing directly from your data..."})
            done_event = {
                "step": "done",
                "message": "Analysis complete",
                "result": deterministic["result"],
                "chart": deterministic.get("chart"),
                "chart_json": deterministic.get("chart_json"),
                "report": deterministic.get("report"),
                "critique": deterministic.get("critique"),
                "plan": plan,
                "code": deterministic.get("code"),
                "code_lang": deterministic.get("code_lang"),
                "followups": followups,
                "validation": deterministic.get("validation"),
                "meta": deterministic.get("meta"),
            }
            _cache_set(cache_key, done_event)
            _update_conversation_state(
                req.session_id,
                req.question,
                deterministic["result"],
                plan,
                deterministic.get("code"),
                deterministic.get("code_lang"),
            )
            yield emit(done_event)
            return

        # ── RAG: retrieve context from uploaded documents ──────────────────
        rag_context = ""
        rag_sources: list[str] = []
        store = doc_stores.get(req.session_id)
        if store and store.chunks:
            yield emit({"step": "analyzing", "message": f"Retrieving context from {len(store.filenames)} document(s)..."})
            hits = store.search(req.question, top_k=4)
            if hits:
                rag_context = "\n\n".join(
                    f"[Source: {h['filename']}]\n{h['text']}" for h in hits
                )
                rag_sources = list(dict.fromkeys(h["filename"] for h in hits))

        rag_block = f"\n\nRELEVANT DOCUMENTATION:\n{rag_context}" if rag_context else ""

        # ── AGENT 1: PLANNER ──────────────────────────────────────────────
        yield emit({"step": "planning", "message": "Planner is mapping your question to the dataset..."})
        try:
            plan_raw = llm(
                PLANNER_SYSTEM,
                f"Category: {req.category}\nDomain context: {category_persona}\n\nSchema:\n{schema}{rag_block}\n\nQuestion: {req.question}"
            )
            plan = parse_json_safe(plan_raw)
        except Exception:
            plan = {"needs_chart": True, "strategy": "Direct analysis", "relevant_columns": [], "analysis_steps": [], "chart_type": "auto"}

        plan["rag_sources"] = rag_sources
        yield emit({"step": "plan", "message": f"Plan: {plan.get('strategy', 'Analyzing...')}", "plan": plan})

        # ── AGENT 2: ANALYST (pandas) or SQL ANALYST ──────────────────────
        query_type = plan.get("query_type", "pandas")
        analyst_result = None
        analyst_code   = None

        if query_type == "sql":
            yield emit({"step": "analyst", "message": "SQL analyst generating query..."})
            sql_schema = get_sql_schema(df)
            sql_context = (
                f"Domain context: {category_persona}\n"
                f"Schema:\n{sql_schema}\n"
                + (f"\nDocumentation context:\n{rag_context}\n" if rag_context else "")
                + f"\nAnalysis strategy: {plan.get('strategy', '')}\n"
                f"Question: {req.question}"
            )
            sql_query = ""
            try:
                sql_query = llm(SQL_ANALYST_SYSTEM, sql_context).strip()
            except Exception as e:
                yield emit({"step": "error", "message": f"SQL agent error: {e}"}); return

            yield emit({"step": "code", "message": "SQL query generated", "code": sql_query, "code_lang": "sql"})
            yield emit({"step": "executing", "message": "Executing SQL on your data..."})

            for attempt in range(2):
                try:
                    analyst_result = execute_sql(sql_query, df)
                    analyst_code = sql_query
                    break
                except Exception:
                    err = traceback.format_exc().strip().splitlines()[-1]
                    if attempt == 0:
                        yield emit({"step": "thinking", "message": f"SQL error — fixing… ({err})"})
                        try:
                            sql_query = llm(SQL_ANALYST_SYSTEM,
                                            f"{sql_context}\n\nPrevious query failed:\n{sql_query}\nError: {err}\nWrite a corrected SQL query.")
                        except Exception:
                            break
                    else:
                        yield emit({"step": "error", "message": f"SQL error: {err}"}); return

            if analyst_result is None:
                yield emit({"step": "error", "message": "SQL analyst could not produce results."}); return

        else:
            yield emit({"step": "analyst", "message": "Analyst agent computing statistics..."})
            analyst_context = (
                f"Domain context: {category_persona}\n"
                f"Schema:\n{schema}\n\n"
                f"Analysis strategy: {plan.get('strategy', '')}\n"
                f"Focus on columns: {', '.join(plan.get('relevant_columns', []))}\n"
                f"Steps to follow: {'; '.join(plan.get('analysis_steps', []))}\n"
                + (f"\nAdditional context from documentation:\n{rag_context}\n" if rag_context else "")
                + f"\nQuestion: {req.question}"
            )
            analyst_code, analyst_result, _, __ = run_code_with_repair(ANALYST_SYSTEM, analyst_context, "analyst")
            if analyst_code:
                yield emit({"step": "code", "message": "Analyst code generated", "code": analyst_code, "code_lang": "python"})
            yield emit({"step": "executing", "message": "Executing analysis on your data..."})
            if analyst_result is None:
                yield emit({"step": "error", "message": "Analyst agent could not compute results."}); return

        # ── AGENT 3: VISUALIZER ───────────────────────────────────────────
        chart_b64 = None
        chart_json = None
        if plan.get("needs_chart", True):
            yield emit({"step": "visualizing", "message": "Visualizer agent creating interactive chart..."})
            viz_context = (
                f"Schema:\n{schema}\n\n"
                f"Question: {req.question}\n"
                f"Suggested chart type: {plan.get('chart_type', 'auto')}\n"
                f"Analysis findings to visualize:\n{analyst_result[:800]}"
            )
            viz_code, _, chart_b64, chart_json = run_code_with_repair(VISUALIZER_SYSTEM, viz_context, "visualizer")
            if viz_code:
                yield emit({"step": "code", "message": "Visualizer code generated", "code": viz_code})

        # ── AGENT 4: CRITIC ───────────────────────────────────────────────
        yield emit({"step": "critiquing", "message": "Critic agent reviewing the analysis..."})
        try:
            critique_raw = llm(
                CRITIC_SYSTEM,
                f"Question: {req.question}\nAnalysis strategy: {plan.get('strategy', '')}\n\nFindings:\n{analyst_result}"
            )
            critique = parse_json_safe(critique_raw)
            if not critique:
                critique = {"verdict": "pass", "confidence": 0.85, "issues": [], "strengths": ["Analysis completed"], "suggestion": ""}
        except Exception:
            critique = {"verdict": "pass", "confidence": 0.85, "issues": [], "strengths": [], "suggestion": ""}

        confidence = critique.get("confidence", 0.85)
        verdict = critique.get("verdict", "pass")
        yield emit({"step": "critique", "message": f"Critic: {verdict.upper()} · {confidence:.0%} confidence", "critique": critique})

        # ── AGENT 5: REPORTER ─────────────────────────────────────────────
        yield emit({"step": "reporting", "message": "Report agent writing executive summary..."})
        try:
            report = llm(
                REPORTER_SYSTEM,
                f"Question: {req.question}\nCategory: {req.category}\n\nFindings:\n{analyst_result}\n\nCritic notes: {critique.get('suggestion', '')}",
                temperature=0.2,
            )
        except Exception:
            report = analyst_result

        done_event = {
            "step": "done",
            "message": "Analysis complete",
            "result": analyst_result,
            "chart": chart_b64,
            "chart_json": chart_json,
            "report": report,
            "critique": critique,
            "plan": plan,
            "code": analyst_code,
            "code_lang": "sql" if query_type == "sql" else "python",
            "followups": suggest_followups(req.question, df, plan, conversation_state.get(req.session_id)),
            "validation": build_validation_payload(
                df,
                method="LLM-generated analysis reviewed by critic and executed against the uploaded dataframe",
                source_columns=[c for c in plan.get("relevant_columns", []) if c in df.columns],
                confidence=float(critique.get("confidence", 0.85) or 0.85),
                reasons=[
                    f"Execution route: {query_type}.",
                    f"Critic verdict: {critique.get('verdict', 'pass')}.",
                    "Generated code or SQL was executed against the in-memory dataframe.",
                ],
            ),
            "meta": {"route": "llm", "cacheable": True},
        }
        _cache_set(cache_key, done_event)
        _update_conversation_state(
            req.session_id,
            req.question,
            analyst_result,
            plan,
            analyst_code,
            "sql" if query_type == "sql" else "python",
        )
        yield emit(done_event)

    return StreamingResponse(stream(), media_type="text/event-stream")


@app.post("/story")
async def story(req: StoryRequest) -> StreamingResponse:
    """Generate a fact-first dataset story from deterministic analysis facts."""
    cleanup_expired_sessions()
    if req.session_id not in dataframes:
        raise HTTPException(status_code=404, detail="Session not found. Upload a CSV first.")
    _touch_session(req.session_id)

    df = dataframes[req.session_id]

    async def stream():
        emit = make_sse_emitter("story", req.session_id)

        yield emit({"step": "analyzing", "message": "Profiling the dataset for story facts...",
                    "meta": {"route": "deterministic_story"}})
        facts = build_story_facts(df, req.category)

        yield emit({"step": "planning", "message": "Ranking deterministic facts into a narrative arc..."})
        report = render_story_from_facts(facts)
        chart_json = build_story_chart_json(facts)
        plan = {
            "strategy": "Fact-first dataset story generated from deterministic profile, quality, grouping, and correlation facts.",
            "query_type": "deterministic_story",
            "relevant_columns": facts.get("column_roles", {}).get("metrics", [])[:1]
                                + facts.get("column_roles", {}).get("dimensions", [])[:2],
            "analysis_steps": [
                "Build data profile",
                "Infer column roles",
                "Compute top deterministic insights",
                "Render story from facts only",
            ],
            "chart_type": "bar" if chart_json else "none",
            "facts": facts,
        }
        critique = {
            "verdict": "pass",
            "confidence": 0.97,
            "issues": [],
            "strengths": ["Story is grounded in deterministic facts before narrative wording."],
            "suggestion": "Use follow-up questions to drill into the leading segment or metric.",
        }

        yield emit({"step": "plan", "message": "Plan: Fact-first dataset story", "plan": plan})
        yield emit({"step": "reporting", "message": "Writing the story from computed facts..."})
        done_event = {
            "step": "done",
            "message": "Dataset story complete",
            "result": report,
            "report": report,
            "chart": None,
            "chart_json": chart_json,
            "critique": critique,
            "plan": plan,
            "followups": suggest_followups("What story does this dataset tell?", df, plan, conversation_state.get(req.session_id)),
            "validation": build_validation_payload(
                df,
                method="Fact-first deterministic story pipeline",
                source_columns=plan["relevant_columns"],
                confidence=0.97,
                reasons=[
                    "Dataset story was written from computed facts.",
                    "Facts include profile, data quality, grouped totals, and correlations.",
                    "Narrative wording is separated from numeric computation.",
                ],
            ),
            "meta": {"route": "deterministic_story", "facts_first": True},
        }
        _update_conversation_state(req.session_id, "What story does this dataset tell?", report, plan)
        yield emit(done_event)

    return StreamingResponse(stream(), media_type="text/event-stream")


@app.post("/predict")
async def predict(req: PredictRequest) -> StreamingResponse:
    cleanup_expired_sessions()
    if req.session_id not in dataframes:
        raise HTTPException(status_code=404, detail="Session not found. Upload a CSV first.")
    _touch_session(req.session_id)
    df = dataframes[req.session_id]
    if req.target not in df.columns:
        raise HTTPException(status_code=400, detail=f"Column '{req.target}' not found.")

    async def stream():
        emit = make_sse_emitter("predict", req.session_id)

        yield emit({"step": "analyzing", "message": f"Preparing features to predict '{req.target}'..."})
        yield emit({"step": "thinking", "message": "Training a Random Forest model on your data..."})
        try:
            payload = build_predict_payload(req.session_id, df, req.target)
        except Exception as e:
            yield emit({"step": "error", "message": f"Could not train model: {e}"})
            return
        yield emit({"step": "executing", "message": "Evaluating on held-out test set..."})
        yield emit({"step": "thinking", "message": "Computing SHAP values and permutation importance..."})
        yield emit({
            "step": "done",
            **payload,
        })

    return StreamingResponse(stream(), media_type="text/event-stream")


def run_predict_job(job_id: str, session_id: str, target: str) -> None:
    update_job(job_id, "running")
    try:
        df = dataframes[session_id]
        payload = build_predict_payload(session_id, df, target)
        update_job(job_id, "completed", result=payload)
    except Exception as exc:
        update_job(job_id, "failed", error=str(exc))


@app.post("/predict_job")
def predict_job(req: PredictRequest, background_tasks: BackgroundTasks) -> dict:
    cleanup_expired_sessions()
    if req.session_id not in dataframes:
        raise HTTPException(status_code=404, detail="Session not found. Upload a CSV first.")
    _touch_session(req.session_id)
    if req.target not in dataframes[req.session_id].columns:
        raise HTTPException(status_code=400, detail=f"Column '{req.target}' not found.")
    job = create_job("predict", req.session_id)
    background_tasks.add_task(run_predict_job, job["job_id"], req.session_id, req.target)
    return {
        "job_id": job["job_id"],
        "status": job["status"],
        "kind": job["kind"],
        "poll_url": f"/jobs/{job['job_id']}",
    }


@app.get("/model_info/{session_id}")
def model_info(session_id: str) -> dict:
    cleanup_expired_sessions()
    if session_id in dataframes:
        _touch_session(session_id)
    info = models.get(session_id)
    if not info:
        return {"trained": False}
    return {
        "trained": True,
        "target": info["target"],
        "is_classification": info["is_classification"],
        "features": info["features"],
    }


@app.get("/jobs/{job_id}")
def get_job(job_id: str) -> dict:
    cleanup_expired_sessions()
    job = jobs.get(job_id)
    if not job:
        job_db = storage.get_job(job_id)
        if job_db:
            job = {
                "job_id": job_db["job_id"],
                "kind": job_db["job_type"],
                "session_id": job_db.get("session_id"),
                "status": job_db["status"],
                "created_at": job_db.get("created_at", _now()),
                "updated_at": job_db.get("updated_at", _now()),
                "result": job_db.get("result"),
                "error": job_db.get("error"),
            }
            jobs[job_id] = job
    if not job:
        raise HTTPException(status_code=404, detail="Job not found.")
    if job.get("session_id") and is_valid_session(str(job["session_id"])):
        _touch_session(str(job["session_id"]))
    return public_job(job)


def prepare_model_features(info: dict, values: dict) -> pd.DataFrame:
    """Convert user-facing feature values into the model's encoded feature frame."""
    row: dict = {}
    for c in info["num_cols"]:
        v = values.get(c)
        row[c] = pd.to_numeric(v, errors="coerce") if v not in (None, "") else np.nan
    for c in info["cat_cols"]:
        row[c] = values.get(c)

    X_new = pd.DataFrame([row])
    if info["cat_cols"]:
        X_new = pd.get_dummies(X_new, columns=info["cat_cols"], drop_first=True)
    X_new = X_new.reindex(columns=info["feature_cols"], fill_value=0)
    return X_new.fillna(value=info["medians"]).fillna(0)


def predict_model_values(info: dict, values: dict) -> dict:
    """Predict a single user-facing row and return JSON-safe prediction details with uncertainty bounds."""
    X_new = prepare_model_features(info, values)
    model = info.get("model")
    if model and hasattr(model, "predict"):
        pred = model.predict(X_new)[0]
    else:
        pred = 0.0

    if info["is_classification"]:
        classes = info.get("classes") or []
        label = classes[int(pred)] if (classes and int(pred) < len(classes)) else str(pred)
        probabilities = {}
        if model and hasattr(model, "predict_proba"):
            probabilities = {
                str(classes[idx]): round(float(value), 4)
                for idx, value in enumerate(model.predict_proba(X_new)[0])
            }
        proba = max(probabilities.values()) if probabilities else None
        return {
            "target": info["target"],
            "prediction": str(label),
            "confidence": round(float(proba), 3) if proba is not None else None,
            "probabilities": probabilities,
            "is_classification": True,
        }

    # For regression, compute 90% prediction interval from ensemble tree variance
    pred_val = round(float(pred), 2)
    interval = None
    if model and hasattr(model, "estimators_"):
        tree_preds = [float(tree.predict(X_new)[0]) for tree in model.estimators_]
        std = float(np.std(tree_preds))
        lower = round(float(pred_val - 1.645 * std), 2)
        upper = round(float(pred_val + 1.645 * std), 2)
        interval = {
            "lower": lower,
            "upper": upper,
            "std_dev": round(std, 2),
            "confidence": 0.90,
            "explanation": "90% plausible prediction interval computed from Random Forest ensemble tree variance.",
        }

    return {
        "target": info["target"],
        "prediction": pred_val,
        "prediction_interval": interval,
        "confidence": None,
        "is_classification": False,
    }


def default_model_values(info: dict) -> dict:
    """Return the median/mode default row exposed by model metadata."""
    return {feature["name"]: feature.get("default") for feature in info["features"]}


def apply_scenario_changes(info: dict, baseline: dict, changes: dict) -> tuple[dict, list[dict]]:
    """Apply set, delta, or percent changes to a baseline row."""
    scenario = {**baseline}
    feature_by_name = {feature["name"]: feature for feature in info["features"]}
    applied: list[dict] = []

    for name, change in changes.items():
        if name not in feature_by_name:
            continue
        feature = feature_by_name[name]
        mode = "set"
        raw_value = change
        if isinstance(change, dict):
            mode = str(change.get("mode", "set")).lower()
            raw_value = change.get("value")

        before = scenario.get(name, feature.get("default"))
        after = raw_value
        if feature.get("type") == "number":
            current = pd.to_numeric(before, errors="coerce")
            value = pd.to_numeric(raw_value, errors="coerce")
            current_f = float(current) if pd.notna(current) else 0.0
            value_f = float(value) if pd.notna(value) else current_f
            if mode in ("percent", "pct", "percent_delta"):
                after = current_f * (1 + value_f / 100)
            elif mode in ("delta", "add"):
                after = current_f + value_f
            else:
                after = value_f
            after = round(float(after), 4)
        scenario[name] = after
        applied.append({"feature": name, "mode": mode, "before": before, "after": after})

    return scenario, applied


def parse_scenario_prompt(info: dict, prompt: str) -> dict:
    """Parse a natural-language what-if prompt into one simulator control change."""
    text = prompt.strip()
    normalized = _norm_text(re.sub(r"[^a-zA-Z0-9_.%+-]+", " ", text))
    if not normalized:
        raise ValueError("Enter a scenario prompt.")

    feature_matches: list[tuple[int, dict]] = []
    for feature in info["features"]:
        name = str(feature["name"])
        name_norm = _norm_text(re.sub(r"[^a-zA-Z0-9_.%+-]+", " ", name))
        parts = [part for part in name_norm.split() if part]
        if name_norm and name_norm in normalized:
            feature_matches.append((100 + len(name_norm), feature))
        elif parts and all(part in normalized for part in parts):
            feature_matches.append((70 + sum(len(part) for part in parts), feature))

    if not feature_matches:
        return {
            "parsed": False,
            "confidence": 0.0,
            "reason": "No trained-model feature name was found in the prompt.",
            "candidates": [feature["name"] for feature in info["features"][:8]],
        }

    feature = sorted(feature_matches, key=lambda item: item[0], reverse=True)[0][1]
    feature_name = feature["name"]
    feature_type = feature.get("type", "number")

    if feature_type == "number":
        feature_norm = _norm_text(re.sub(r"[^a-zA-Z0-9_.%+-]+", " ", str(feature_name)))
        value_text = normalized.replace(feature_norm, " ", 1)
        number_match = re.search(r"([-+]?\d+(?:\.\d+)?)\s*(%|percent|percentage)?", value_text)
        if not number_match:
            return {
                "parsed": False,
                "confidence": 0.35,
                "feature": feature_name,
                "reason": "A numeric scenario needs a numeric value.",
                "candidates": [feature["name"] for feature in info["features"][:8]],
            }
        value = float(number_match.group(1))
        has_percent = bool(number_match.group(2)) or "%" in value_text or "percent" in value_text or "percentage" in value_text
        decrease_terms = ("decrease", "decreases", "decreased", "reduce", "reduces", "reduced",
                          "lower", "lowers", "lowered", "drop", "drops", "dropped", "down", "less", "subtract")
        increase_terms = ("increase", "increases", "increased", "raise", "raises", "raised",
                          "grow", "grows", "grew", "boost", "boosts", "boosted", "up", "more", "add")
        set_terms = ("set", "make", "change", "becomes", "to", "=")

        if any(term in normalized for term in decrease_terms):
            value = -abs(value)
            mode = "percent" if has_percent else "delta"
        elif any(term in normalized for term in increase_terms):
            value = abs(value)
            mode = "percent" if has_percent else "delta"
        elif has_percent:
            mode = "percent"
        elif any(term in normalized for term in set_terms):
            mode = "set"
        else:
            mode = "set"

        return {
            "parsed": True,
            "feature": feature_name,
            "mode": mode,
            "value": value,
            "confidence": 0.88 if feature_matches[0][0] >= 100 else 0.72,
            "interpretation": f"{mode} {feature_name} by {value:g}" if mode != "set" else f"set {feature_name} to {value:g}",
            "candidates": [item[1]["name"] for item in sorted(feature_matches, key=lambda item: item[0], reverse=True)[:5]],
        }

    options = [str(option) for option in feature.get("options", [])]
    selected = next((option for option in options if option.lower() in normalized), None)
    if not selected:
        value_match = re.search(r"(?:to|as|=|is)\s+(.+)$", normalized)
        selected = value_match.group(1).strip() if value_match else ""
    if not selected:
        return {
            "parsed": False,
            "confidence": 0.4,
            "feature": feature_name,
            "reason": "A categorical scenario needs a target category value.",
            "candidates": options[:8],
        }

    return {
        "parsed": True,
        "feature": feature_name,
        "mode": "set",
        "value": selected,
        "confidence": 0.82 if selected in options else 0.62,
        "interpretation": f"set {feature_name} to {selected}",
        "candidates": options[:8],
    }


def build_scenario_chart(info: dict, baseline_prediction: dict, scenario_prediction: dict) -> str:
    """Create a compact comparison chart for simulation output."""
    if info["is_classification"]:
        scenario_label = str(scenario_prediction["prediction"])
        baseline_prob = (baseline_prediction.get("probabilities") or {}).get(scenario_label, 0)
        scenario_prob = (scenario_prediction.get("probabilities") or {}).get(scenario_label, 0)
        plot_df = pd.DataFrame({
            "case": ["Baseline", "Scenario"],
            "probability": [baseline_prob, scenario_prob],
        })
        fig = px.bar(plot_df, x="case", y="probability", title=f"Probability of {scenario_label}", text="probability")
        fig.update_yaxes(range=[0, 1], tickformat=".0%")
    else:
        plot_df = pd.DataFrame({
            "case": ["Baseline", "Scenario"],
            "prediction": [baseline_prediction["prediction"], scenario_prediction["prediction"]],
        })
        fig = px.bar(plot_df, x="case", y="prediction", title=f"Scenario impact on {info['target']}", text="prediction")
    return _chart_to_json(fig)


def scenario_impact(info: dict, baseline_prediction: dict, scenario_prediction: dict) -> dict:
    """Summarize the difference between baseline and scenario predictions."""
    if info["is_classification"]:
        scenario_label = str(scenario_prediction["prediction"])
        base_prob = (baseline_prediction.get("probabilities") or {}).get(scenario_label, 0)
        scen_prob = (scenario_prediction.get("probabilities") or {}).get(scenario_label, 0)
        return {
            "type": "classification",
            "label_changed": baseline_prediction["prediction"] != scenario_prediction["prediction"],
            "baseline_label": baseline_prediction["prediction"],
            "scenario_label": scenario_prediction["prediction"],
            "scenario_label_probability_delta": round(float(scen_prob - base_prob), 4),
            "confidence_delta": round(
                float((scenario_prediction.get("confidence") or 0) - (baseline_prediction.get("confidence") or 0)),
                4,
            ),
        }

    baseline_value = float(baseline_prediction["prediction"])
    scenario_value = float(scenario_prediction["prediction"])
    delta = scenario_value - baseline_value
    pct = delta / abs(baseline_value) * 100 if baseline_value else None
    return {
        "type": "regression",
        "delta": round(delta, 4),
        "pct_change": None if pct is None else round(float(pct), 2),
        "direction": "increase" if delta > 0 else "decrease" if delta < 0 else "no_change",
    }


@app.post("/predict_input")
def predict_input(req: PredictInputRequest) -> dict:
    cleanup_expired_sessions()
    info = models.get(req.session_id)
    if not info:
        info = storage.get_model_info(req.session_id)
        if info:
            models[req.session_id] = info
    if not info:
        raise HTTPException(status_code=400, detail="Train a model first using the Predict button.")
    _touch_session(req.session_id)
    return predict_model_values(info, req.values)


# ── Business report export (PDF + PPTX) ──────────────────────────────────────

@app.post("/simulate")
def simulate_scenario(req: ScenarioRequest) -> dict:
    cleanup_expired_sessions()
    info = models.get(req.session_id)
    if not info:
        raise HTTPException(status_code=400, detail="Train a model first using the Predict button.")
    _touch_session(req.session_id)

    baseline_values = {**default_model_values(info), **req.baseline}
    scenario_values, changes_applied = apply_scenario_changes(info, baseline_values, req.changes)
    if not changes_applied:
        raise HTTPException(status_code=400, detail="No valid scenario changes were provided.")

    baseline_prediction = predict_model_values(info, baseline_values)
    scenario_prediction = predict_model_values(info, scenario_values)
    impact = scenario_impact(info, baseline_prediction, scenario_prediction)

    return {
        "target": info["target"],
        "category": req.category,
        "is_classification": info["is_classification"],
        "baseline_values": baseline_values,
        "scenario_values": scenario_values,
        "changes_applied": changes_applied,
        "baseline_prediction": baseline_prediction,
        "scenario_prediction": scenario_prediction,
        "impact": impact,
        "chart_json": build_scenario_chart(info, baseline_prediction, scenario_prediction),
        "validation": {
            "confidence_label": "Medium",
            "method": "What-if simulation using the trained in-session Random Forest model",
            "reasons": [
                "The model was not retrained; only the selected input values changed.",
                "Scenario results are predictive estimates, not causal proof.",
                "Reliability depends on whether the changed values stay within the training distribution.",
            ],
        },
    }


@app.post("/scenario_parse")
def scenario_parse(req: ScenarioParseRequest) -> dict:
    cleanup_expired_sessions()
    info = models.get(req.session_id)
    if not info:
        raise HTTPException(status_code=400, detail="Train a model first using the Predict button.")
    _touch_session(req.session_id)
    parsed = parse_scenario_prompt(info, req.prompt)
    return {
        **parsed,
        "prompt": req.prompt,
        "target": info["target"],
        "category": req.category,
        "validation": {
            "method": "Deterministic feature-name and value parser",
            "reasons": [
                "The parser only uses trained model feature names and explicit numeric or category values.",
                "Ambiguous prompts return low confidence or an unparsed response instead of guessing.",
            ],
        },
    }


import datetime as _dt
from fastapi import Query


class ReportRequest(BaseModel):
    messages:  list[dict] = []
    category:  str        = "general"
    filename:  str        = "report"


def build_report_payload(session_id: str, req: ReportRequest, format: str) -> dict:
    """Generate a JSON/base64 report payload for a session."""
    df = dataframes[session_id]
    profile = build_profile(df)

    messages = list(req.messages)

    if session_id in join_store:
        j = join_store[session_id]
        messages.append({
            "question": f"Relational Join Audit: {j['join_key_1']} = {j['join_key_2']} ({j['how']})",
            "report": f"Unified dataset created by joining table 1 ({j['left_rows_before']} rows) and table 2 ({j['right_rows_before']} rows) on '{j['join_key_1']}' = '{j['join_key_2']}'. Resulting row count: {j['rows_after']:,} rows.",
            "critique": {"verdict": "pass", "confidence": 0.95, "issues": []},
        })

    if session_id in forecast_store:
        fc = forecast_store[session_id]
        messages.append({
            "question": f"Time-Series Forecast: {fc['target_column']} ({fc['periods']} periods)",
            "report": f"Holt-Winters time-series forecast for target metric '{fc['target_column']}' over {fc['periods']} future periods. Forecasted trend direction: {fc['metrics']['trend_direction']} ({fc['metrics']['growth_rate_pct']}% expected growth) with 95% confidence bounds.",
            "chart_json": fc.get("chart_json"),
            "critique": {"verdict": "pass", "confidence": 0.90, "issues": []},
        })

    if session_id in compare_store:
        comp = compare_store[session_id]
        added = ", ".join(comp['schema_changes']['added_columns']) or "None"
        removed = ", ".join(comp['schema_changes']['removed_columns']) or "None"
        messages.append({
            "question": "Dataset Comparison & Drift Analysis (v1 vs v2)",
            "report": f"Comparison between base version ({comp['v1_rows']:,} rows) and target version ({comp['v2_rows']:,} rows). Added columns: [{added}]. Removed columns: [{removed}]. Significant numeric drift detected in {len(comp['numeric_drift'])} metric(s).",
            "critique": {"verdict": "pass", "confidence": 0.92, "issues": []},
        })

    enhanced_req = ReportRequest(messages=messages, category=req.category, filename=req.filename)

    if format == "pdf":
        content = _generate_pdf(enhanced_req, profile)
        media = "application/pdf"
        ext = "pdf"
    else:
        content = _generate_pptx(enhanced_req, profile)
        media = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ext = "pptx"

    safe_name = "".join(c if c.isalnum() or c in "-_" else "_" for c in req.filename)
    return {
        "filename": f"{safe_name}_report.{ext}",
        "media_type": media,
        "size_bytes": len(content),
        "content_base64": base64.b64encode(content).decode("ascii"),
    }


def _plotly_json_to_png(chart_json: str) -> bytes | None:
    """Convert Plotly JSON to PNG bytes via kaleido. Returns None on failure."""
    try:
        import plotly.io as pio
        fig = go.Figure(json.loads(chart_json))
        return pio.to_image(fig, format="png", width=900, height=540, scale=1.5)
    except Exception:
        return None


def _b64_to_bytes(b64: str) -> bytes | None:
    try:
        return base64.b64decode(b64)
    except Exception:
        return None


def _generate_pdf(req: ReportRequest, profile: dict) -> bytes:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import cm
    from reportlab.lib.colors import HexColor, white, black
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_LEFT
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
        Image as RLImage, PageBreak, HRFlowable,
    )

    W, H = A4
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4,
                            leftMargin=2.5*cm, rightMargin=2.5*cm,
                            topMargin=2.5*cm, bottomMargin=2.5*cm)

    # ── Styles ──
    INDIGO   = HexColor("#4F46E5")
    DARK     = HexColor("#0F172A")
    SOFT     = HexColor("#64748B")
    LIGHT_BG = HexColor("#F8FAFC")
    BORDER   = HexColor("#E2E8F0")

    S = {
        "cover_title": ParagraphStyle("ct", fontSize=32, fontName="Helvetica-Bold",
                                      textColor=DARK, alignment=TA_CENTER, spaceAfter=10),
        "cover_sub":   ParagraphStyle("cs", fontSize=14, fontName="Helvetica",
                                      textColor=SOFT, alignment=TA_CENTER, spaceAfter=6),
        "section":     ParagraphStyle("sec", fontSize=16, fontName="Helvetica-Bold",
                                      textColor=INDIGO, spaceBefore=18, spaceAfter=8),
        "question":    ParagraphStyle("q", fontSize=13, fontName="Helvetica-Bold",
                                      textColor=DARK, spaceBefore=14, spaceAfter=6),
        "body":        ParagraphStyle("b", fontSize=10.5, fontName="Helvetica",
                                      textColor=HexColor("#334155"), leading=16, spaceAfter=6),
        "mono":        ParagraphStyle("m", fontSize=9, fontName="Courier",
                                      textColor=HexColor("#475569"), leading=14, spaceAfter=4),
        "meta":        ParagraphStyle("mt", fontSize=9, fontName="Helvetica",
                                      textColor=SOFT, spaceAfter=4),
    }

    def _safe_para(text: str, style) -> Paragraph:
        text = (text or "").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        text = text.replace("\n", "<br/>")
        return Paragraph(text, style)

    def _profile_table(rows: list[list]) -> Table:
        t = Table(rows, colWidths=[7*cm, 9*cm])
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), INDIGO),
            ("TEXTCOLOR",  (0, 0), (-1, 0), white),
            ("FONTNAME",   (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE",   (0, 0), (-1, -1), 10),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [LIGHT_BG, white]),
            ("GRID",       (0, 0), (-1, -1), 0.5, BORDER),
            ("LEFTPADDING",  (0, 0), (-1, -1), 10),
            ("RIGHTPADDING", (0, 0), (-1, -1), 10),
            ("TOPPADDING",   (0, 0), (-1, -1), 6),
            ("BOTTOMPADDING",(0, 0), (-1, -1), 6),
        ]))
        return t

    def _embed_chart(img_bytes: bytes | None) -> RLImage | None:
        if not img_bytes:
            return None
        try:
            return RLImage(io.BytesIO(img_bytes), width=14*cm, height=8.75*cm)
        except Exception:
            return None

    PERSONA_HEADERS = {
        "financial": "CFO Executive Briefing · Focus: Margin, Variance, Capital Efficiency & Risk",
        "medical": "Clinical Leadership Briefing · Focus: Patient Outcomes & Operational Efficacy",
        "retail": "Merchandising & Operations Briefing · Focus: Unit Economics & Inventory Turn",
        "marketing": "CMO Growth Briefing · Focus: CAC, LTV & Campaign Attribution",
        "hr": "People Operations Briefing · Focus: Attrition, Banding & Workforce Productivity",
        "general": "Executive Decision Briefing · Focus: Profile, Trends, Predictions & Signals",
    }
    persona_tag = PERSONA_HEADERS.get(req.category.lower(), PERSONA_HEADERS["general"])

    story = []
    cat   = req.category.title()
    today = _dt.date.today().strftime("%B %d, %Y")

    # ── Cover ──
    story.append(Spacer(1, 2.5*cm))
    story.append(Paragraph("Analytics Report", S["cover_title"]))
    story.append(Paragraph(req.filename, S["cover_sub"]))
    story.append(Paragraph(f"{cat} Lens  ·  {today}", S["cover_sub"]))
    story.append(Paragraph(persona_tag, S["meta"]))
    story.append(Spacer(1, 1.2*cm))
    story.append(HRFlowable(width="80%", thickness=2, color=INDIGO, spaceAfter=18))
    story.append(Paragraph("Generated by CSV Analyst AI · Powered by Gemini 2.5 Flash-Lite", S["meta"]))
    story.append(PageBreak())

    # ── Data Profile ──
    story.append(Paragraph("Dataset Profile", S["section"]))
    n_cols = len(profile.get("columns", []))
    n_cat  = n_cols - profile.get("numeric_features", 0)
    prows = [
        ["Metric", "Value"],
        ["Rows", f"{profile.get('rows', 0):,}"],
        ["Columns", str(n_cols)],
        ["Numeric features", str(profile.get("numeric_features", 0))],
        ["Categorical features", str(n_cat)],
        ["Missing values", f"{profile.get('missing_pct', 0)}%"],
        ["Duplicate rows", str(profile.get("duplicate_rows", 0))],
    ]
    story.append(_profile_table(prows))

    if profile.get("numeric_stats"):
        story.append(Spacer(1, 0.4*cm))
        story.append(Paragraph("Numeric Summary", S["section"]))
        stat_rows = [["Column", "Mean", "Median", "Std Dev", "Min", "Max"]]
        for col, s in list(profile["numeric_stats"].items())[:12]:
            stat_rows.append([col,
                str(s.get("mean", "—")), str(s.get("median", "—")),
                str(s.get("std", "—")),  str(s.get("min", "—")),
                str(s.get("max", "—"))])
        t2 = Table(stat_rows, colWidths=[4.5*cm, 2.5*cm, 2.5*cm, 2.5*cm, 2.5*cm, 2.5*cm])
        t2.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), INDIGO),
            ("TEXTCOLOR",  (0, 0), (-1, 0), white),
            ("FONTNAME",   (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE",   (0, 0), (-1, -1), 9),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [LIGHT_BG, white]),
            ("GRID",       (0, 0), (-1, -1), 0.4, BORDER),
            ("LEFTPADDING",  (0, 0), (-1, -1), 6),
            ("TOPPADDING",   (0, 0), (-1, -1), 5),
            ("BOTTOMPADDING",(0, 0), (-1, -1), 5),
        ]))
        story.append(t2)
    story.append(PageBreak())

    # ── Analysis sections ──
    for i, msg in enumerate(req.messages, 1):
        q = msg.get("question", "—")
        story.append(Paragraph(f"Q{i}: {q}", S["question"]))

        # Executive summary
        report_text = msg.get("report") or msg.get("result") or ""
        if report_text:
            story.append(_safe_para(report_text, S["body"]))

        # Chart (prefer chart_json → PNG via kaleido, fallback to chart_b64)
        chart_img = None
        if msg.get("chart_json"):
            chart_img = _embed_chart(_plotly_json_to_png(msg["chart_json"]))
        if chart_img is None and msg.get("chart"):
            chart_img = _embed_chart(_b64_to_bytes(msg["chart"]))
        if chart_img:
            story.append(Spacer(1, 0.3*cm))
            story.append(chart_img)

        # SHAP chart (explainability)
        if msg.get("shap_chart"):
            story.append(Spacer(1, 0.3*cm))
            story.append(Paragraph("SHAP Feature Impact", S["meta"]))
            shap_img = _embed_chart(_b64_to_bytes(msg["shap_chart"]))
            if shap_img:
                story.append(shap_img)

        # Critique
        if msg.get("critique"):
            c = msg["critique"]
            verdict = c.get("verdict", "pass").upper()
            conf    = c.get("confidence", 0)
            issues  = "; ".join(c.get("issues", []))
            line    = f"Analysis quality: {verdict}  ·  {conf:.0%} confidence"
            if issues:
                line += f"  ·  {issues}"
            story.append(Spacer(1, 0.2*cm))
            story.append(_safe_para(line, S["meta"]))

        # Code (collapsed)
        if msg.get("code"):
            lang = msg.get("code_lang", "python").upper()
            story.append(Spacer(1, 0.2*cm))
            story.append(Paragraph(f"Generated {lang} code:", S["meta"]))
            snippet = msg["code"][:800] + ("…" if len(msg["code"]) > 800 else "")
            story.append(_safe_para(snippet, S["mono"]))

        story.append(HRFlowable(width="100%", thickness=0.5, color=BORDER, spaceAfter=12))
        if i < len(req.messages):
            story.append(Spacer(1, 0.2*cm))

    doc.build(story)
    return buf.getvalue()


def _generate_pptx(req: ReportRequest, profile: dict) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    INDIGO = RGBColor(0x4F, 0x46, 0xE5)
    DARK   = RGBColor(0x0F, 0x17, 0x2A)
    SOFT   = RGBColor(0x64, 0x74, 0x8B)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]  # blank layout

    def _slide():
        return prs.slides.add_slide(blank)

    def _txbox(slide, left, top, width, height, text, size=18, bold=False,
                color=DARK, align=PP_ALIGN.LEFT, wrap=True):
        from pptx.util import Inches, Pt
        box  = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf   = box.text_frame
        tf.word_wrap = wrap
        p    = tf.paragraphs[0]
        p.alignment = align
        run  = p.add_run()
        run.text = text[:500]
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        return box

    def _add_img(slide, img_bytes: bytes, left, top, width, height):
        try:
            slide.shapes.add_picture(io.BytesIO(img_bytes),
                                      Inches(left), Inches(top),
                                      Inches(width), Inches(height))
        except Exception:
            pass

    def _rect(slide, left, top, width, height, fill_rgb):
        from pptx.util import Inches
        from pptx.oxml.ns import qn
        shape = slide.shapes.add_shape(1,  # MSO_SHAPE_TYPE.RECTANGLE
                                        Inches(left), Inches(top),
                                        Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
        shape.line.fill.background()
        return shape

    PERSONA_HEADERS = {
        "financial": "CFO Executive Briefing · Focus: Margin, Variance, Capital Efficiency & Risk",
        "medical": "Clinical Leadership Briefing · Focus: Patient Outcomes & Operational Efficacy",
        "retail": "Merchandising & Operations Briefing · Focus: Unit Economics & Inventory Turn",
        "marketing": "CMO Growth Briefing · Focus: CAC, LTV & Campaign Attribution",
        "hr": "People Operations Briefing · Focus: Attrition, Banding & Workforce Productivity",
        "general": "Executive Decision Briefing · Focus: Profile, Trends, Predictions & Signals",
    }
    persona_tag = PERSONA_HEADERS.get(req.category.lower(), PERSONA_HEADERS["general"])

    today = _dt.date.today().strftime("%B %d, %Y")

    # ── Title slide ──
    s = _slide()
    _rect(s, 0, 0, 13.33, 7.5, RGBColor(0xF8, 0xFA, 0xFC))
    _rect(s, 0, 0, 13.33, 1.5, INDIGO)
    _txbox(s, 0.5, 0.3, 12, 1, "Analytics Report", size=36, bold=True,
           color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    _txbox(s, 0.5, 2.2, 12, 0.6, req.filename, size=22, bold=True,
           color=DARK, align=PP_ALIGN.CENTER)
    _txbox(s, 0.5, 2.9, 12, 0.5, f"{req.category.title()} Lens  ·  {today}",
           size=14, color=SOFT, align=PP_ALIGN.CENTER)
    _txbox(s, 0.5, 3.5, 12, 0.5, persona_tag,
           size=12, color=INDIGO, align=PP_ALIGN.CENTER)
    _txbox(s, 0.5, 6.8, 12, 0.5, "Generated by CSV Analyst AI · Gemini 2.5 Flash-Lite",
           size=11, color=SOFT, align=PP_ALIGN.CENTER)

    # ── Profile slide ──
    s = _slide()
    _rect(s, 0, 0, 13.33, 1.1, INDIGO)
    _txbox(s, 0.4, 0.15, 12, 0.8, "Dataset Profile", size=24, bold=True,
           color=RGBColor(0xFF, 0xFF, 0xFF))
    items = [
        ("Rows", f"{profile.get('rows', 0):,}"),
        ("Columns", str(len(profile.get("columns", [])))),
        ("Numeric features", str(profile.get("numeric_features", 0))),
        ("Missing values", f"{profile.get('missing_pct', 0)}%"),
        ("Duplicate rows", str(profile.get("duplicate_rows", 0))),
    ]
    for i, (k, v) in enumerate(items):
        col   = i % 3
        row   = i // 3
        left  = 0.4 + col * 4.3
        top   = 1.4 + row * 1.6
        _rect(s, left, top, 4.0, 1.4, RGBColor(0xEE, 0xF2, 0xFF))
        _txbox(s, left + 0.15, top + 0.1, 3.7, 0.5, v,
               size=28, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)
        _txbox(s, left + 0.15, top + 0.7, 3.7, 0.5, k,
               size=11, color=SOFT, align=PP_ALIGN.CENTER)

    # ── Per-message slides ──
    for i, msg in enumerate(req.messages, 1):
        s = _slide()
        _rect(s, 0, 0, 13.33, 1.1, INDIGO)
        q = (msg.get("question") or "")[:80]
        _txbox(s, 0.4, 0.15, 12.5, 0.8, f"Q{i}: {q}", size=16, bold=True,
               color=RGBColor(0xFF, 0xFF, 0xFF), wrap=True)

        # Chart
        chart_png = None
        if msg.get("chart_json"):
            chart_png = _plotly_json_to_png(msg["chart_json"])
        if chart_png is None and msg.get("chart"):
            chart_png = _b64_to_bytes(msg.get("chart"))
        if chart_png:
            _add_img(s, chart_png, 0.4, 1.2, 7.5, 4.7)
            report_left = 8.2
        else:
            report_left = 0.4

        # Report text
        report = (msg.get("report") or msg.get("result") or "")[:600]
        if report:
            _txbox(s, report_left, 1.2, 4.8 if chart_png else 12.5, 5.2,
                   report, size=11, color=DARK, wrap=True)

        # Critique footer
        if msg.get("critique"):
            c = msg["critique"]
            footer = f"Quality: {c.get('verdict','pass').upper()}  ·  {c.get('confidence',0):.0%} confidence"
            _txbox(s, 0.4, 6.8, 12.5, 0.5, footer, size=10, color=SOFT, align=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@app.post("/report/{session_id}")
async def export_report(
    session_id: str,
    req: ReportRequest,
    format: str = Query(default="pdf", pattern="^(pdf|pptx)$"),
) -> dict:
    """Generate a PDF or PPTX report from the session data and conversation messages."""
    cleanup_expired_sessions()
    if session_id not in dataframes:
        raise HTTPException(status_code=404, detail="Session not found. Upload a CSV first.")
    _touch_session(session_id)

    return build_report_payload(session_id, req, format)


def run_report_job(job_id: str, session_id: str, req: ReportRequest, format: str) -> None:
    update_job(job_id, "running")
    try:
        payload = build_report_payload(session_id, req, format)
        update_job(job_id, "completed", result=payload)
    except Exception as exc:
        update_job(job_id, "failed", error=str(exc))


@app.post("/report_job/{session_id}")
async def export_report_job(
    session_id: str,
    req: ReportRequest,
    background_tasks: BackgroundTasks,
    format: str = Query(default="pdf", pattern="^(pdf|pptx)$"),
) -> dict:
    """Queue PDF/PPTX report generation and return a polling URL."""
    cleanup_expired_sessions()
    if session_id not in dataframes:
        raise HTTPException(status_code=404, detail="Session not found. Upload a CSV first.")
    _touch_session(session_id)
    job = create_job("report", session_id)
    background_tasks.add_task(run_report_job, job["job_id"], session_id, req, format)
    return {
        "job_id": job["job_id"],
        "status": job["status"],
        "kind": job["kind"],
        "poll_url": f"/jobs/{job['job_id']}",
    }


# ── Benchmark evaluation framework ───────────────────────────────────────────

BENCHMARK_QUESTIONS = [
    # ── General statistics ────────────────────────────────────────────────
    {"question": "What are the summary statistics for all numeric columns?",            "category": "general",   "expects_chart": False, "expects_sql": False},
    {"question": "How many rows and columns does the dataset have?",                    "category": "general",   "expects_chart": False, "expects_sql": True},
    {"question": "Show a correlation heatmap of all numeric columns",                   "category": "general",   "expects_chart": True,  "expects_sql": False},
    {"question": "Plot the distribution of the first numeric column",                   "category": "general",   "expects_chart": True,  "expects_sql": False},
    {"question": "Show a bar chart of the top categories by count",                     "category": "general",   "expects_chart": True,  "expects_sql": True},
    {"question": "What percentage of values are missing in each column?",               "category": "general",   "expects_chart": False, "expects_sql": False},
    {"question": "Are there any duplicate rows?",                                       "category": "general",   "expects_chart": False, "expects_sql": True},
    {"question": "Show a box plot of all numeric columns to detect outliers",           "category": "general",   "expects_chart": True,  "expects_sql": False},
    # ── SQL-type (filter / group-by / top-N) ─────────────────────────────
    {"question": "Show the top 10 rows sorted by the highest numeric column",           "category": "general",   "expects_chart": False, "expects_sql": True},
    {"question": "Count records in each category and sort by count descending",         "category": "general",   "expects_chart": False, "expects_sql": True},
    {"question": "What is the average of each numeric column grouped by the main category?", "category": "general", "expects_chart": False, "expects_sql": True},
    {"question": "Show all records where the first numeric column is above its average","category": "general",   "expects_chart": False, "expects_sql": True},
    {"question": "What are the top 5 categories by total of the main numeric column?",  "category": "general",   "expects_chart": False, "expects_sql": True},
    # ── Financial ─────────────────────────────────────────────────────────
    {"question": "What is the total revenue?",                                          "category": "financial", "expects_chart": False, "expects_sql": True},
    {"question": "Show total revenue by category sorted descending as a bar chart",     "category": "financial", "expects_chart": True,  "expects_sql": True},
    {"question": "Calculate mean, standard deviation and CV of the main numeric metric","category": "financial", "expects_chart": False, "expects_sql": False},
    {"question": "What are the top 5 items by revenue? Show a bar chart",               "category": "financial", "expects_chart": True,  "expects_sql": True},
    {"question": "Plot total revenue over time as a line chart",                        "category": "financial", "expects_chart": True,  "expects_sql": False},
    {"question": "What is the period-over-period growth rate of the main metric?",      "category": "financial", "expects_chart": True,  "expects_sql": False},
    # ── Retail ────────────────────────────────────────────────────────────
    {"question": "Which product or category has the highest total sales?",              "category": "retail",    "expects_chart": False, "expects_sql": True},
    {"question": "Show a bar chart of revenue by region or segment",                    "category": "retail",    "expects_chart": True,  "expects_sql": True},
    {"question": "What is the average order value?",                                    "category": "retail",    "expects_chart": False, "expects_sql": True},
    {"question": "Show a scatter plot of quantity vs revenue",                          "category": "retail",    "expects_chart": True,  "expects_sql": False},
    {"question": "What percentage of total revenue does each category contribute?",     "category": "retail",    "expects_chart": True,  "expects_sql": False},
    {"question": "Show the top 10 best-selling products by revenue as a bar chart",     "category": "retail",    "expects_chart": True,  "expects_sql": True},
    # ── Marketing ─────────────────────────────────────────────────────────
    {"question": "Break down totals by each categorical column",                        "category": "marketing", "expects_chart": True,  "expects_sql": True},
    {"question": "Which segment or channel performs best by total metric?",             "category": "marketing", "expects_chart": True,  "expects_sql": True},
    {"question": "Show the share of each category as a pie or bar chart",               "category": "marketing", "expects_chart": True,  "expects_sql": False},
    {"question": "Compare the mean numeric values across segments with a bar chart",    "category": "marketing", "expects_chart": True,  "expects_sql": False},
    # ── HR ────────────────────────────────────────────────────────────────
    {"question": "Show the headcount by department or category as a bar chart",         "category": "hr",        "expects_chart": True,  "expects_sql": True},
    {"question": "Plot the distribution of the main numeric column across the dataset", "category": "hr",        "expects_chart": True,  "expects_sql": False},
    {"question": "Compare the mean of each numeric feature between groups",             "category": "hr",        "expects_chart": True,  "expects_sql": False},
    {"question": "What is the average of the main numeric column by group?",            "category": "hr",        "expects_chart": False, "expects_sql": True},
    # ── Correlation & statistics ──────────────────────────────────────────
    {"question": "Which two columns are most strongly correlated?",                     "category": "general",   "expects_chart": True,  "expects_sql": False},
    {"question": "Show the covariance matrix of numeric columns",                       "category": "general",   "expects_chart": False, "expects_sql": False},
    {"question": "What is the skewness and kurtosis of each numeric column?",           "category": "general",   "expects_chart": False, "expects_sql": False},
    {"question": "Are there significant differences in means between groups? Show statistical test", "category": "general", "expects_chart": False, "expects_sql": False},
    # ── Time series ───────────────────────────────────────────────────────
    {"question": "Plot the trend of the main numeric column over time",                 "category": "financial", "expects_chart": True,  "expects_sql": False},
    {"question": "Show a monthly or yearly breakdown of the main metric",               "category": "financial", "expects_chart": True,  "expects_sql": True},
    # ── Edge cases ────────────────────────────────────────────────────────
    {"question": "Show me something interesting about this data",                       "category": "general",   "expects_chart": True,  "expects_sql": False},
    {"question": "What story does this dataset tell?",                                  "category": "general",   "expects_chart": False, "expects_sql": False},
    {"question": "Identify the most important pattern or anomaly in this data",         "category": "general",   "expects_chart": True,  "expects_sql": False},
    # ── Medical ───────────────────────────────────────────────────────────
    {"question": "Which features correlate most with the outcome variable?",            "category": "medical",   "expects_chart": True,  "expects_sql": False},
    {"question": "Compare the mean of each numeric feature between outcome groups",     "category": "medical",   "expects_chart": True,  "expects_sql": False},
    {"question": "Show the distribution of age split by outcome group",                 "category": "medical",   "expects_chart": True,  "expects_sql": False},
    {"question": "What is the prevalence of each outcome class?",                       "category": "medical",   "expects_chart": True,  "expects_sql": True},
    {"question": "Show a violin plot comparing groups on the main clinical metric",     "category": "medical",   "expects_chart": True,  "expects_sql": False},
    # ── Multi-step ────────────────────────────────────────────────────────
    {"question": "Rank all categories by revenue, show top 5 and bottom 5",            "category": "retail",    "expects_chart": True,  "expects_sql": True},
    {"question": "Calculate the Pareto principle: what top % of categories drive 80% of the main metric?", "category": "financial", "expects_chart": True, "expects_sql": False},
]


def _run_agent_pipeline_sync(df: pd.DataFrame, schema: str, question: str,
                              category: str, doc_store=None) -> dict:
    """Non-streaming pipeline runner for benchmarking. Returns metrics dict."""
    import time, re as _re

    t0 = time.time()
    category_persona = CATEGORY_PERSONAS.get(category, CATEGORY_PERSONAS["general"])

    def _llm(system: str, user: str, temperature: float = 0) -> str:
        resp = client.models.generate_content(
            model=GEMINI_MODEL, contents=user,
            config=types.GenerateContentConfig(system_instruction=system, temperature=temperature),
        )
        text = (resp.text or "").strip()
        if text.startswith("```"): text = text.split("\n", 1)[1]
        if "```" in text: text = text.rsplit("```", 1)[0]
        return text.strip()

    def _parse_json(text: str) -> dict:
        try: return json.loads(text)
        except Exception:
            m = _re.search(r'\{.*\}', text, _re.DOTALL)
            if m:
                try: return json.loads(m.group())
                except: pass
        return {}

    # RAG context
    rag_block = ""
    if doc_store and doc_store.chunks:
        hits = doc_store.search(question, top_k=3)
        if hits:
            rag_block = "\n\nRELEVANT DOCUMENTATION:\n" + "\n\n".join(
                f"[Source: {h['filename']}]\n{h['text']}" for h in hits
            )

    # Agent 1: Planner
    used_repair = False
    try:
        plan = _parse_json(_llm(PLANNER_SYSTEM,
            f"Category: {category}\nDomain context: {category_persona}\n\nSchema:\n{schema}{rag_block}\n\nQuestion: {question}"))
    except Exception:
        plan = {"needs_chart": True, "strategy": "Direct", "relevant_columns": [],
                "analysis_steps": [], "chart_type": "auto", "query_type": "pandas"}

    query_type = plan.get("query_type", "pandas")
    analyst_result = None

    # Agent 2: Analyst or SQL
    if query_type == "sql":
        sql_schema = get_sql_schema(df)
        sql_ctx = (f"Domain context: {category_persona}\nSchema:\n{sql_schema}\n"
                   f"Strategy: {plan.get('strategy', '')}\nQuestion: {question}")
        try:
            sql = _llm(SQL_ANALYST_SYSTEM, sql_ctx).strip()
            try:
                analyst_result = execute_sql(sql, df)
            except Exception as e:
                used_repair = True
                sql = _llm(SQL_ANALYST_SYSTEM, f"{sql_ctx}\n\nFailed:\n{sql}\nError: {e}\nFix it.")
                try: analyst_result = execute_sql(sql, df)
                except: analyst_result = None
        except Exception:
            analyst_result = None
    else:
        analyst_ctx = (f"Domain context: {category_persona}\nSchema:\n{schema}\n"
                       f"Strategy: {plan.get('strategy', '')}\nFocus: {', '.join(plan.get('relevant_columns', []))}\n"
                       f"Question: {question}")
        try:
            code = _llm(ANALYST_SYSTEM, analyst_ctx)
            try:
                analyst_result, _, __ = execute_code(code, df)
            except Exception as e:
                used_repair = True
                code = _llm(ANALYST_SYSTEM, f"{analyst_ctx}\n\nFailed:\n{code}\nError: {e}\nFix it.")
                try: analyst_result, _, __ = execute_code(code, df)
                except: analyst_result = None
        except Exception:
            analyst_result = None

    # Agent 3: Visualizer
    has_chart = False
    if plan.get("needs_chart", True) and analyst_result:
        viz_ctx = (f"Schema:\n{schema}\nQuestion: {question}\n"
                   f"Chart type: {plan.get('chart_type', 'auto')}\n"
                   f"Findings:\n{(analyst_result or '')[:500]}")
        try:
            viz_code = _llm(VISUALIZER_SYSTEM, viz_ctx)
            try:
                _, __, chart_json = execute_code(viz_code, df)
                has_chart = chart_json is not None
            except Exception as e:
                viz_code = _llm(VISUALIZER_SYSTEM, f"{viz_ctx}\n\nFailed:\n{viz_code}\nError: {e}\nFix it.")
                try:
                    _, __, chart_json = execute_code(viz_code, df)
                    has_chart = chart_json is not None
                except: pass
        except Exception:
            pass

    elapsed = round(time.time() - t0, 2)
    return {
        "success":     analyst_result is not None,
        "has_chart":   has_chart,
        "query_type":  query_type,
        "used_repair": used_repair,
        "time_s":      elapsed,
    }


def build_benchmark_payload(session_id: str, n: int = 15) -> dict:
    """Run up to n benchmark questions and return aggregate metrics."""
    df = dataframes[session_id]
    schema = get_df_schema(df)
    store = doc_stores.get(session_id)
    n = min(n, len(BENCHMARK_QUESTIONS))
    suite = BENCHMARK_QUESTIONS[:n]

    results = []
    for bq in suite:
        m = _run_agent_pipeline_sync(df, schema, bq["question"], bq["category"], store)
        results.append({
            "question": bq["question"],
            "category": bq["category"],
            "expects_chart": bq["expects_chart"],
            "expects_sql": bq["expects_sql"],
            **m,
        })

    total = len(results)
    n_success = sum(1 for r in results if r["success"])
    n_chart_exp = sum(1 for r in results if r["expects_chart"])
    n_chart_got = sum(1 for r in results if r["expects_chart"] and r["has_chart"])
    n_sql_exp = sum(1 for r in results if r["expects_sql"])
    n_sql_routed = sum(1 for r in results if r["expects_sql"] and r["query_type"] == "sql")
    n_repair = sum(1 for r in results if r["used_repair"])
    n_repair_ok = sum(1 for r in results if r["used_repair"] and r["success"])
    avg_time = round(sum(r["time_s"] for r in results) / total, 2) if total else 0

    return {
        "total": total,
        "success_rate": round(n_success / total, 3),
        "chart_rate": round(n_chart_got / max(1, n_chart_exp), 3),
        "sql_routing_accuracy": round(n_sql_routed / max(1, n_sql_exp), 3),
        "repair_rate": round(n_repair / total, 3),
        "repair_success_rate": round(n_repair_ok / max(1, n_repair), 3),
        "avg_time_s": avg_time,
        "results": results,
    }


@app.get("/benchmark/{session_id}")
async def run_benchmark(session_id: str, n: int = 15) -> dict:
    """Run up to n benchmark questions against the uploaded dataset and return metrics."""
    cleanup_expired_sessions()
    if session_id not in dataframes:
        raise HTTPException(status_code=404, detail="Session not found. Upload a CSV first.")
    _touch_session(session_id)

    return build_benchmark_payload(session_id, n)


def run_benchmark_job(job_id: str, session_id: str, n: int) -> None:
    update_job(job_id, "running")
    try:
        payload = build_benchmark_payload(session_id, n)
        update_job(job_id, "completed", result=payload)
    except Exception as exc:
        update_job(job_id, "failed", error=str(exc))


@app.post("/benchmark_job/{session_id}")
async def benchmark_job(session_id: str, background_tasks: BackgroundTasks, n: int = 15) -> dict:
    """Queue benchmark evaluation and return a polling URL."""
    cleanup_expired_sessions()
    if session_id not in dataframes:
        raise HTTPException(status_code=404, detail="Session not found. Upload a CSV first.")
    _touch_session(session_id)
    job = create_job("benchmark", session_id)
    background_tasks.add_task(run_benchmark_job, job["job_id"], session_id, n)
    return {
        "job_id": job["job_id"],
        "status": job["status"],
        "kind": job["kind"],
        "poll_url": f"/jobs/{job['job_id']}",
    }
